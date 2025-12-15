"""
现代化PPT编辑器 - 仿PowerPoint界面
UI风格：参考PowerPoint的现代布局
- 顶部：红色标题栏 + 双行工具栏
- 左侧：页面缩略图导航
- 中间：主编辑画布
- 右侧：属性面板
- 底部：红色状态栏
"""

import tkinter as tk
from tkinter import filedialog, messagebox, ttk, colorchooser
from PIL import Image, ImageTk, ImageDraw, ImageFont
import json
import os
import threading
import logging
import cv2
import numpy as np
import tempfile
import copy
from datetime import datetime

# PDF支持 - 使用PyMuPDF，更简单，不需要Poppler
try:
    import fitz  # PyMuPDF
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("提示: 安装 PyMuPDF 可支持PDF导入")
    print("      pip install PyMuPDF")


logging.getLogger("ppocr").setLevel(logging.WARNING)

# 获取程序运行目录
def get_base_dir():
    import sys
    if getattr(sys, 'frozen', False):
        # 打包后的exe运行目录
        return os.path.dirname(sys.executable)
    else:
        # 开发环境
        return os.path.dirname(os.path.abspath(__file__))

# 配置文件路径
CONFIG_FILE = os.path.join(get_base_dir(), "ppt_editor_config.json")

def load_config():
    """加载配置"""
    default_config = {
        "model_dir": os.path.join(get_base_dir(), ".paddlex", "official_models")
    }
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                config = json.load(f)
                # 合并默认配置
                for key in default_config:
                    if key not in config:
                        config[key] = default_config[key]
                return config
        except:
            pass
    return default_config

def save_config(config):
    """保存配置"""
    try:
        with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
            json.dump(config, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"保存配置失败: {e}")

try:
    from paddleocr import PaddleOCR
except ImportError:
    print("请先安装 paddleocr: pip install paddleocr paddlepaddle")
    exit()

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# === 全局配色（仿PowerPoint） ===
COLOR_THEME = "#B7472A"           # PowerPoint红色主题
COLOR_THEME_HOVER = "#C85A3F"     # 悬停色
COLOR_RIBBON_BG = "#F5F5F5"       # Ribbon工具栏背景
COLOR_RIBBON_ROW2 = "#E8E8E8"     # 第二行背景
COLOR_CANVAS_BG = "#E0E0E0"       # 画布背景
COLOR_SIDEBAR_BG = "#FAFAFA"      # 侧边栏背景
COLOR_WHITE = "#FFFFFF"
COLOR_TEXT = "#333333"
COLOR_BLUE = "#1976D2"
COLOR_GREEN = "#43A047"
COLOR_ORANGE = "#FB8C00"
COLOR_PURPLE = "#8E24AA"
COLOR_RED = "#E53935"
COLOR_GRAY = "#607D8B"
FONT_FAMILY = "微软雅黑"


def Px(pixels):
    """像素转EMU单位"""
    return Emu(int(pixels) * 9525)


class TextBox:
    """文本框数据类"""
    def __init__(self, x, y, width, height):
        self.x = x
        self.y = y
        self.width = width
        self.height = height
        self.text = ""
        self.font_size = 16
        self.font_name = "微软雅黑"
        self.font_color = "#000000"
        self.bold = False
        self.italic = False
        self.align = "left"

    def to_dict(self):
        return {
            "x": self.x, "y": self.y, "width": self.width, "height": self.height,
            "text": self.text, "font_size": self.font_size, "font_name": self.font_name,
            "font_color": self.font_color, "bold": self.bold, "italic": self.italic,
            "align": self.align
        }

    @staticmethod
    def from_dict(data):
        box = TextBox(data["x"], data["y"], data["width"], data["height"])
        box.text = data.get("text", "")
        box.font_size = data.get("font_size", 16)
        box.font_name = data.get("font_name", "微软雅黑")
        box.font_color = data.get("font_color", "#000000")
        box.bold = data.get("bold", False)
        box.italic = data.get("italic", False)
        box.align = data.get("align", "left")
        return box

    def copy(self):
        """复制文本框"""
        return copy.deepcopy(self)


class ModernPPTEditor:
    def __init__(self, root):
        self.root = root
        self.root.title("PPT编辑器专业版 - 增强版")
        self.root.geometry("1500x900")
        self.root.configure(bg=COLOR_RIBBON_BG)

        # 加载配置
        self.config = load_config()

        # 多页支持
        self.pages = []
        self.current_page_index = 0

        # 当前页数据
        self.original_img_path = None
        self.clean_bg_path = None
        self.original_image = None
        self.display_image = None
        self.tk_image = None
        self.scale = 1.0

        # 文本框
        self.text_boxes = []
        self.selected_box_index = -1
        self.selected_boxes = []

        # 预览模式
        self.current_preview_mode = "original"
        self.ppt_preview_image = None

        # 撤销/重做
        self.history = []
        self.history_index = -1
        self.max_history = 50

        # 绘制状态
        self.is_drawing = False
        self.draw_start_x = 0
        self.draw_start_y = 0
        self.temp_rect_id = None
        self.is_dragging = False
        self.drag_start_x = 0
        self.drag_start_y = 0
        self.is_resizing = False
        self.resize_handle = None
        self.is_selecting = False  # 框选模式
        self.select_start_x = 0
        self.select_start_y = 0

        # 绘制模式
        self.draw_mode = True

        # OCR模型
        self.ocr = None

        # 缩略图
        self.thumbnail_images = []
        # 复制粘贴支持
        self.clipboard_boxes = []

        # 自动保存
        self.autosave_timer = None
        self.project_file_path = None
        self.has_unsaved_changes = False

        # 创建自动保存目录
        AUTOSAVE_DIR = os.path.join(get_base_dir(), "autosave")
        os.makedirs(AUTOSAVE_DIR, exist_ok=True)
        self.autosave_dir = AUTOSAVE_DIR


        # 创建界面
        self.create_ui()

        # 绑定快捷键
        self.bind_shortcuts()

        # 后台加载OCR
        threading.Thread(target=self.init_ocr, daemon=True).start()

        # 启动自动保存
        if self.config.get("autosave_enabled", True):
            self.start_autosave()

        # 窗口关闭事件
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)


    def _imread_unicode(self, filepath):
        """
        安全读取包含中文路径的图片
        解决OpenCV无法读取中文路径的问题
        """
        try:
            # 使用numpy读取文件，然后解码为图片
            img_array = np.fromfile(filepath, dtype=np.uint8)
            img = cv2.imdecode(img_array, cv2.IMREAD_COLOR)
            return img
        except Exception as e:
            print(f"读取图片失败: {filepath}, 错误: {e}")
            return None

    def init_ocr(self):
        """后台初始化OCR模型 - 优先使用配置的模型目录"""
        self.update_status("正在加载OCR模型...")
        try:
            # 从配置获取模型目录和设备
            local_model_dir = self.config.get("model_dir", os.path.join(get_base_dir(), ".paddlex", "official_models"))
            device = self.config.get("ocr_device", "cpu")  # 从配置读取设备（cpu 或 gpu）

            # 所有模型路径
            det_model = os.path.join(local_model_dir, "PP-OCRv5_server_det")
            rec_model = os.path.join(local_model_dir, "PP-OCRv5_server_rec")
            doc_ori_model = os.path.join(local_model_dir, "PP-LCNet_x1_0_doc_ori")
            textline_ori_model = os.path.join(local_model_dir, "PP-LCNet_x1_0_textline_ori")
            unwarp_model = os.path.join(local_model_dir, "UVDoc")

            # 检查核心模型是否存在
            if os.path.exists(det_model) and os.path.exists(rec_model):
                # 构建参数
                ocr_params = {
                    "lang": "ch",
                    "text_detection_model_dir": det_model,
                    "text_recognition_model_dir": rec_model,
                    "device": device  # 使用配置的设备
                }

                # 可选模型：存在则使用本地，否则禁用
                if os.path.exists(doc_ori_model):
                    ocr_params["doc_orientation_classify_model_dir"] = doc_ori_model
                if os.path.exists(textline_ori_model):
                    ocr_params["textline_orientation_model_dir"] = textline_ori_model
                if os.path.exists(unwarp_model):
                    ocr_params["doc_unwarping_model_dir"] = unwarp_model

                self.ocr = PaddleOCR(**ocr_params)
                device_name = "GPU" if device == "gpu" else "CPU"
                self.update_status(f"OCR模型加载完成（本地模型，{device_name}）")
            else:
                # 回退到默认（自动下载）
                self.ocr = PaddleOCR(lang="ch")
                self.update_status("OCR模型加载完成")
        except Exception as e:
            self.update_status(f"OCR加载失败: {e}")
            import traceback
            traceback.print_exc()

    def create_ui(self):
        """创建界面"""
        # === 顶部标题栏 ===
        self.create_title_bar()

        # === 工具栏 ===
        self.create_toolbar()

        # === 主内容区 ===
        self.main_container = tk.Frame(self.root, bg=COLOR_CANVAS_BG)
        self.main_container.pack(fill=tk.BOTH, expand=True)

        # 左侧：页面缩略图
        self.create_thumbnail_panel()

        # 中间：主编辑区
        self.create_canvas_area()

        # 右侧：属性面板
        self.create_property_panel()

        # === 底部状态栏 ===
        self.create_status_bar()

    def create_title_bar(self):
        """创建顶部标题栏 - PowerPoint红色风格"""
        title_bar = tk.Frame(self.root, bg=COLOR_THEME, height=32)
        title_bar.pack(fill=tk.X, side=tk.TOP)
        title_bar.pack_propagate(False)

        # 左侧标题
        title_label = tk.Label(title_bar, text="PPT编辑器专业版",
                              bg=COLOR_THEME, fg="white",
                              font=(FONT_FAMILY, 11, "bold"))
        title_label.pack(side=tk.LEFT, padx=15)

        # 右侧页码信息
        self.title_page_label = tk.Label(title_bar, text="第 0/0 页",
                                         bg=COLOR_THEME, fg="white",
                                         font=(FONT_FAMILY, 10))
        self.title_page_label.pack(side=tk.RIGHT, padx=15)

        # 自动保存状态指示器
        self.autosave_indicator = tk.Label(title_bar, text="●",
                                          bg=COLOR_THEME, fg="#4CAF50",
                                          font=(FONT_FAMILY, 16))
        self.autosave_indicator.pack(side=tk.RIGHT, padx=5)

    def create_toolbar(self):
        """创建顶部工具栏 - 双行紧凑版（仿PowerPoint Ribbon）"""
        toolbar = tk.Frame(self.root, bg=COLOR_RIBBON_BG, relief=tk.FLAT)
        toolbar.pack(fill=tk.X, side=tk.TOP)

        # 底部边框线
        border_line = tk.Frame(toolbar, bg="#ddd", height=1)
        border_line.pack(fill=tk.X, side=tk.BOTTOM)

        # === 第一行：文件和OCR操作 ===
        row1 = tk.Frame(toolbar, bg=COLOR_RIBBON_BG)
        row1.pack(fill=tk.X, padx=10, pady=(6, 2))

        # 文件组
        tk.Label(row1, text="文件:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row1, "导入图片", self.load_multiple_images, COLOR_GREEN)
        self.create_tool_btn(row1, "导入背景", self.load_multiple_backgrounds, COLOR_BLUE)
        if PDF_SUPPORT:
            self.create_tool_btn(row1, "导入PDF", self.import_pdf, "#D32F2F")

        self.create_tool_btn(row1, "保存项目", self.save_project, COLOR_GRAY)
        self.create_tool_btn(row1, "打开项目", self.load_project, COLOR_GRAY)

        self.create_separator(row1)

        # 检测组
        tk.Label(row1, text="检测:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row1, "当前页", self.auto_detect_text_regions, COLOR_ORANGE)
        self.create_tool_btn(row1, "全部页", self.auto_detect_all_pages, "#EF6C00")

        self.create_separator(row1)

        # 识别组
        tk.Label(row1, text="识别:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row1, "当前页", self.ocr_all_boxes, COLOR_PURPLE)
        self.create_tool_btn(row1, "全部页", self.ocr_all_pages, "#6A1B9A")

        self.create_separator(row1)

        # 自动字号组
        tk.Label(row1, text="自动字号:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row1, "当前页", self.auto_font_size_all, "#00ACC1")
        self.create_tool_btn(row1, "全部页", self.auto_font_size_all_pages, "#00838F")


        self.create_separator(row1)

        # 导出组
        tk.Label(row1, text="导出:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row1, "生成PPT", self.generate_multi_page_ppt, COLOR_RED)
        self.create_tool_btn(row1, "导出PDF", self.export_as_pdf, "#C62828")
        self.create_tool_btn(row1, "导出图片", self.export_as_images, "#F57C00")

        # 右侧设置按钮
        settings_btn = tk.Button(row1, text="⚙ 设置", command=self.show_settings_dialog,
                                bg="#546E7A", fg="white", font=(FONT_FAMILY, 9),
                                padx=8, pady=2, cursor="hand2", relief=tk.FLAT, bd=0)
        settings_btn.pack(side=tk.RIGHT, padx=5)

        # === 第二行：编辑和预览 ===
        row2 = tk.Frame(toolbar, bg=COLOR_RIBBON_ROW2)
        row2.pack(fill=tk.X, padx=10, pady=(2, 6))

        # 编辑工具
        tk.Label(row2, text="编辑:", bg=COLOR_RIBBON_ROW2, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.draw_mode_var = tk.BooleanVar(value=True)
        self.draw_btn = tk.Button(row2, text="画框模式", command=self.toggle_draw_mode_btn,
                                  bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 9),
                                  padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        self.draw_btn.pack(side=tk.LEFT, padx=2)


        self.create_tool_btn(row2, "复制", self.copy_boxes, "#009688", bg=COLOR_RIBBON_ROW2)
        self.create_tool_btn(row2, "粘贴", self.paste_boxes, "#00ACC1", bg=COLOR_RIBBON_ROW2)
        self.create_tool_btn(row2, "删除框", self.delete_selected_box, COLOR_RED, bg=COLOR_RIBBON_ROW2)
        self.create_tool_btn(row2, "清空全部", self.clear_all_boxes, "#795548", bg=COLOR_RIBBON_ROW2)
        self.create_tool_btn(row2, "撤销", self.undo, "#78909C", bg=COLOR_RIBBON_ROW2)
        self.create_tool_btn(row2, "重做", self.redo, "#78909C", bg=COLOR_RIBBON_ROW2)
        self.create_separator(row2, bg=COLOR_RIBBON_ROW2)

        # 预览模式
        tk.Label(row2, text="预览:", bg=COLOR_RIBBON_ROW2, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.preview_mode_var = tk.StringVar(value="original")
        self.preview_orig_btn = tk.Button(row2, text="原图", command=lambda: self.set_preview_mode("original"),
                                          bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9),
                                          padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        self.preview_orig_btn.pack(side=tk.LEFT, padx=2)

        self.preview_ppt_btn = tk.Button(row2, text="PPT效果", command=lambda: self.set_preview_mode("ppt"),
                                         bg="#757575", fg="white", font=(FONT_FAMILY, 9),
                                         padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        self.preview_ppt_btn.pack(side=tk.LEFT, padx=2)

        self.create_separator(row2, bg=COLOR_RIBBON_ROW2)

        # 缩放控制
        tk.Label(row2, text="视图:", bg=COLOR_RIBBON_ROW2, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row2, "适应窗口", self.fit_image_to_canvas, "#455A64", bg=COLOR_RIBBON_ROW2)
        self.create_tool_btn(row2, "100%", self.zoom_to_100, "#455A64", bg=COLOR_RIBBON_ROW2)

        self.zoom_label = tk.Label(row2, text="100%", bg=COLOR_RIBBON_ROW2, fg="#333",
                                   font=(FONT_FAMILY, 9), padx=10)
        self.zoom_label.pack(side=tk.LEFT)

        # 快捷键提示
        tk.Label(row2, text="Ctrl+滚轮缩放 | 双击编辑 | Ctrl+点击多选 | 拖动调整位置/大小",
                bg=COLOR_RIBBON_ROW2, fg="#999", font=(FONT_FAMILY, 8)).pack(side=tk.LEFT, padx=10)

    def create_tool_btn(self, parent, text, command, color, bg=None):
        """创建工具栏按钮"""
        if bg is None:
            bg = COLOR_RIBBON_BG
        btn = tk.Button(parent, text=text, command=command,
                       bg=color, fg="white", font=(FONT_FAMILY, 9),
                       padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        btn.pack(side=tk.LEFT, padx=2)
        return btn

    def create_separator(self, parent, bg=None):
        """创建分隔线"""
        if bg is None:
            bg = COLOR_RIBBON_BG
        sep_frame = tk.Frame(parent, bg=bg)
        sep_frame.pack(side=tk.LEFT, padx=6)
        sep_line = tk.Frame(sep_frame, bg="#ccc", width=1, height=20)
        sep_line.pack()

    def create_icon_button(self, parent, text, command, color, icon=""):
        """创建图标按钮"""
        btn_text = f"{icon}\n{text}" if icon else text
        btn = tk.Button(parent, text=btn_text, command=command,
                       bg=color, fg="white", font=("微软雅黑", 8),
                       width=5, height=2, cursor="hand2", relief=tk.GROOVE, bd=2)
        btn.pack(side=tk.LEFT, padx=2, pady=2)

        # 悬停效果
        def on_enter(e):
            btn.config(relief=tk.RAISED)
        def on_leave(e):
            btn.config(relief=tk.GROOVE)
        btn.bind("<Enter>", on_enter)
        btn.bind("<Leave>", on_leave)
        return btn

    def toggle_draw_mode_btn(self):
        """切换绘制模式"""
        self.draw_mode = not self.draw_mode
        self.draw_mode_var.set(self.draw_mode)
        if self.draw_mode:
            self.draw_btn.config(bg=COLOR_GREEN, text="画框模式")
            self.canvas.config(cursor="crosshair")
        else:
            self.draw_btn.config(bg="#9E9E9E", text="选择模式")
            self.canvas.config(cursor="")

    def set_preview_mode(self, mode):
        """设置预览模式"""
        self.preview_mode_var.set(mode)
        self.current_preview_mode = mode
        if mode == "original":
            self.preview_orig_btn.config(bg=COLOR_BLUE, fg="white")
            self.preview_ppt_btn.config(bg="#757575", fg="white")
        else:
            self.preview_orig_btn.config(bg="#757575", fg="white")
            self.preview_ppt_btn.config(bg=COLOR_BLUE, fg="white")
        self.refresh_canvas()

    def create_thumbnail_panel(self):
        """创建左侧缩略图面板"""
        self.thumbnail_panel = tk.Frame(self.main_container, bg=COLOR_SIDEBAR_BG, width=180)
        self.thumbnail_panel.pack(side=tk.LEFT, fill=tk.Y)
        self.thumbnail_panel.pack_propagate(False)

        # 标题栏
        title_frame = tk.Frame(self.thumbnail_panel, bg=COLOR_BLUE, height=30)
        title_frame.pack(fill=tk.X)
        title_frame.pack_propagate(False)
        tk.Label(title_frame, text="  页面列表", bg=COLOR_BLUE, fg="white",
                font=(FONT_FAMILY, 10, "bold"), anchor="w").pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # 缩略图容器（可滚动）
        container = tk.Frame(self.thumbnail_panel, bg=COLOR_SIDEBAR_BG)
        container.pack(fill=tk.BOTH, expand=True)

        self.thumbnail_canvas = tk.Canvas(container, bg=COLOR_SIDEBAR_BG, highlightthickness=0, width=160)
        scrollbar = tk.Scrollbar(container, orient=tk.VERTICAL, command=self.thumbnail_canvas.yview)

        self.thumbnail_frame = tk.Frame(self.thumbnail_canvas, bg=COLOR_SIDEBAR_BG)

        self.thumbnail_canvas.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.thumbnail_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        self.thumbnail_window = self.thumbnail_canvas.create_window((0, 0), window=self.thumbnail_frame, anchor=tk.NW)

        self.thumbnail_frame.bind("<Configure>",
            lambda e: self.thumbnail_canvas.configure(scrollregion=self.thumbnail_canvas.bbox("all")))

        # 鼠标滚轮
        self.thumbnail_canvas.bind("<MouseWheel>",
            lambda e: self.thumbnail_canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        # 页面导航按钮
        nav_frame = tk.Frame(self.thumbnail_panel, bg="#f5f5f5", height=40)
        nav_frame.pack(fill=tk.X, side=tk.BOTTOM)
        nav_frame.pack_propagate(False)

        tk.Button(nav_frame, text="上一页", command=self.prev_page,
                 bg="#e0e0e0", font=(FONT_FAMILY, 9), width=6, cursor="hand2",
                 relief=tk.FLAT).pack(side=tk.LEFT, padx=5, pady=5)

        self.page_label = tk.Label(nav_frame, text="0/0", bg="#f5f5f5",
                                   font=(FONT_FAMILY, 10, "bold"))
        self.page_label.pack(side=tk.LEFT, expand=True)

        tk.Button(nav_frame, text="下一页", command=self.next_page,
                 bg="#e0e0e0", font=(FONT_FAMILY, 9), width=6, cursor="hand2",
                 relief=tk.FLAT).pack(side=tk.RIGHT, padx=5, pady=5)

    def create_canvas_area(self):
        """创建中间画布区域"""
        canvas_container = tk.Frame(self.main_container, bg=COLOR_CANVAS_BG)
        canvas_container.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # 画布
        self.canvas = tk.Canvas(canvas_container, bg="#c0c0c0", highlightthickness=0)

        # 滚动条
        v_scroll = tk.Scrollbar(canvas_container, orient=tk.VERTICAL, command=self.canvas.yview)
        h_scroll = tk.Scrollbar(canvas_container, orient=tk.HORIZONTAL, command=self.canvas.xview)

        self.canvas.config(xscrollcommand=h_scroll.set, yscrollcommand=v_scroll.set)

        v_scroll.pack(side=tk.RIGHT, fill=tk.Y)
        h_scroll.pack(side=tk.BOTTOM, fill=tk.X)
        self.canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # 绑定画布事件
        self.canvas.bind("<ButtonPress-1>", self.on_canvas_press)
        self.canvas.bind("<Control-ButtonPress-1>", self.on_canvas_ctrl_click)
        self.canvas.bind("<B1-Motion>", self.on_canvas_drag)
        self.canvas.bind("<ButtonRelease-1>", self.on_canvas_release)
        self.canvas.bind("<Configure>", self.on_canvas_resize)
        self.canvas.bind("<Double-Button-1>", self.on_canvas_double_click)
        self.canvas.bind("<Button-3>", self.on_canvas_right_click)  # 右键菜单

        # Ctrl+滚轮缩放
        self.canvas.bind("<Control-MouseWheel>", self.on_canvas_zoom)
        # 普通滚轮滚动
        self.canvas.bind("<MouseWheel>", self.on_canvas_scroll)

        # 占位提示
        self.placeholder_label = tk.Label(self.canvas,
            text="点击上方「导入图片」按钮开始\n\n支持批量导入多张图片",
            bg="#c0c0c0", fg="#666666", font=(FONT_FAMILY, 14), justify=tk.CENTER)
        self.canvas.create_window(400, 300, window=self.placeholder_label)

    def create_property_panel(self):
        """创建右侧属性面板"""
        self.right_panel = tk.Frame(self.main_container, bg=COLOR_WHITE, width=280)
        self.right_panel.pack(side=tk.RIGHT, fill=tk.Y)
        self.right_panel.pack_propagate(False)

        # 标题
        title_frame = tk.Frame(self.right_panel, bg=COLOR_BLUE, height=30)
        title_frame.pack(fill=tk.X)
        title_frame.pack_propagate(False)
        tk.Label(title_frame, text="  属性设置", bg=COLOR_BLUE, fg="white",
                font=(FONT_FAMILY, 10, "bold"), anchor="w").pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # 可滚动容器
        canvas = tk.Canvas(self.right_panel, bg=COLOR_WHITE, highlightthickness=0)
        scrollbar = tk.Scrollbar(self.right_panel, orient=tk.VERTICAL, command=canvas.yview)

        self.prop_frame = tk.Frame(canvas, bg=COLOR_WHITE)

        canvas.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        canvas_window = canvas.create_window((0, 0), window=self.prop_frame, anchor=tk.NW)

        self.prop_frame.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.bind("<Configure>", lambda e: canvas.itemconfig(canvas_window, width=e.width))
        canvas.bind("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        # === 文本框列表 ===
        self.create_section_header(self.prop_frame, "文本框列表")

        list_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        list_frame.pack(fill=tk.X, padx=10, pady=5)

        self.box_listbox = tk.Listbox(list_frame, height=5, bg="#f5f5f5",
                                       font=(FONT_FAMILY, 9), selectbackground=COLOR_BLUE,
                                       selectforeground="white", relief=tk.FLAT, bd=1)
        self.box_listbox.pack(fill=tk.X)
        self.box_listbox.bind("<<ListboxSelect>>", self.on_listbox_select)

        # === 文本内容 ===
        self.create_section_header(self.prop_frame, "文本内容")

        text_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        text_frame.pack(fill=tk.X, padx=10, pady=5)

        self.text_entry = tk.Text(text_frame, height=3, bg="#f5f5f5",
                                  font=(FONT_FAMILY, 10), relief=tk.FLAT, bd=1, wrap=tk.WORD)
        self.text_entry.pack(fill=tk.X)
        self.text_entry.bind("<KeyRelease>", self.on_text_change)

        # OCR识别按钮
        ocr_btn_frame = tk.Frame(text_frame, bg=COLOR_WHITE)
        ocr_btn_frame.pack(fill=tk.X, pady=5)

        tk.Button(ocr_btn_frame, text="🔍 OCR识别此框", command=self.ocr_single_box,
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 9, "bold"),
                 cursor="hand2", relief=tk.FLAT).pack(fill=tk.X)

        # === 位置和大小 ===
        self.create_section_header(self.prop_frame, "位置和大小")

        pos_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        pos_frame.pack(fill=tk.X, padx=10, pady=5)

        # X, Y
        row1 = tk.Frame(pos_frame, bg=COLOR_WHITE)
        row1.pack(fill=tk.X, pady=2)

        tk.Label(row1, text="X:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9), width=3).pack(side=tk.LEFT)
        self.x_entry = tk.Entry(row1, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.x_entry.pack(side=tk.LEFT, padx=2)
        self.x_entry.bind("<KeyRelease>", self.on_position_change)

        tk.Label(row1, text="Y:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9), width=3).pack(side=tk.LEFT, padx=(10, 0))
        self.y_entry = tk.Entry(row1, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.y_entry.pack(side=tk.LEFT, padx=2)
        self.y_entry.bind("<KeyRelease>", self.on_position_change)

        # 宽, 高
        row2 = tk.Frame(pos_frame, bg=COLOR_WHITE)
        row2.pack(fill=tk.X, pady=2)

        tk.Label(row2, text="宽:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9), width=3).pack(side=tk.LEFT)
        self.w_entry = tk.Entry(row2, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.w_entry.pack(side=tk.LEFT, padx=2)
        self.w_entry.bind("<KeyRelease>", self.on_position_change)

        tk.Label(row2, text="高:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9), width=3).pack(side=tk.LEFT, padx=(10, 0))
        self.h_entry = tk.Entry(row2, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.h_entry.pack(side=tk.LEFT, padx=2)
        self.h_entry.bind("<KeyRelease>", self.on_position_change)

        # === 字体样式 ===
        self.create_section_header(self.prop_frame, "字体样式")

        font_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        font_frame.pack(fill=tk.X, padx=10, pady=5)

        # 字体和字号
        row3 = tk.Frame(font_frame, bg=COLOR_WHITE)
        row3.pack(fill=tk.X, pady=2)

        self.fontname_var = tk.StringVar(value="微软雅黑")
        font_combo = ttk.Combobox(row3, textvariable=self.fontname_var, width=10,
                                  values=["微软雅黑", "宋体", "黑体", "楷体", "仿宋", "Arial"])
        font_combo.pack(side=tk.LEFT, padx=2)
        font_combo.bind("<<ComboboxSelected>>", self.on_font_change)

        self.fontsize_var = tk.StringVar(value="16")
        size_combo = ttk.Combobox(row3, textvariable=self.fontsize_var, width=5,
                                  values=["8", "10", "12", "14", "16", "18", "20", "24", "28", "32", "36", "48", "60", "72", "80", "100", "120", "150", "200"])
        size_combo.pack(side=tk.LEFT, padx=2)
        size_combo.bind("<<ComboboxSelected>>", self.on_font_change)

        # 样式按钮
        row4 = tk.Frame(font_frame, bg=COLOR_WHITE)
        row4.pack(fill=tk.X, pady=5)

        self.bold_var = tk.BooleanVar(value=False)
        self.bold_btn = tk.Button(row4, text="B 加粗", command=self.toggle_bold,
                                  bg="#e0e0e0", font=(FONT_FAMILY, 9),
                                  width=6, cursor="hand2", relief=tk.FLAT)
        self.bold_btn.pack(side=tk.LEFT, padx=2)

        self.italic_var = tk.BooleanVar(value=False)
        self.italic_btn = tk.Button(row4, text="I 斜体", command=self.toggle_italic,
                                    bg="#e0e0e0", font=(FONT_FAMILY, 9),
                                    width=6, cursor="hand2", relief=tk.FLAT)
        self.italic_btn.pack(side=tk.LEFT, padx=2)

        self.color_btn = tk.Button(row4, text="颜色", command=self.choose_color,
                                   bg="#000000", fg="white", width=5, cursor="hand2", relief=tk.FLAT)
        self.color_btn.pack(side=tk.LEFT, padx=2)

        # 自动字号按钮
        tk.Button(row4, text="自动字号", command=self.auto_font_size,
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 8),
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=5)

        # 对齐按钮
        row5 = tk.Frame(font_frame, bg=COLOR_WHITE)
        row5.pack(fill=tk.X, pady=5)

        tk.Label(row5, text="对齐:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.align_var = tk.StringVar(value="left")

        align_btn_frame = tk.Frame(row5, bg=COLOR_WHITE)
        align_btn_frame.pack(side=tk.LEFT, padx=5)

        self.align_left_btn = tk.Button(align_btn_frame, text="左", command=lambda: self.set_align("left"),
                                        bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9), width=3,
                                        cursor="hand2", relief=tk.FLAT)
        self.align_left_btn.pack(side=tk.LEFT, padx=1)

        self.align_center_btn = tk.Button(align_btn_frame, text="中", command=lambda: self.set_align("center"),
                                          bg="#e0e0e0", fg="#333", font=(FONT_FAMILY, 9), width=3,
                                          cursor="hand2", relief=tk.FLAT)
        self.align_center_btn.pack(side=tk.LEFT, padx=1)

        self.align_right_btn = tk.Button(align_btn_frame, text="右", command=lambda: self.set_align("right"),
                                         bg="#e0e0e0", fg="#333", font=(FONT_FAMILY, 9), width=3,
                                         cursor="hand2", relief=tk.FLAT)
        self.align_right_btn.pack(side=tk.LEFT, padx=1)

        # === 批量应用 ===
        self.create_section_header(self.prop_frame, "批量应用")

        batch_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        batch_frame.pack(fill=tk.X, padx=10, pady=5)

        tk.Label(batch_frame, text="Ctrl+点击多选，勾选要应用的属性：",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8), wraplength=220).pack(anchor="w")

        # 勾选项
        check_row1 = tk.Frame(batch_frame, bg=COLOR_WHITE)
        check_row1.pack(fill=tk.X, pady=2)

        self.apply_fontsize_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row1, text="字号", variable=self.apply_fontsize_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.apply_fontname_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row1, text="字体", variable=self.apply_fontname_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.apply_color_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row1, text="颜色", variable=self.apply_color_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        check_row2 = tk.Frame(batch_frame, bg=COLOR_WHITE)
        check_row2.pack(fill=tk.X, pady=2)

        self.apply_bold_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row2, text="加粗", variable=self.apply_bold_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.apply_italic_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row2, text="斜体", variable=self.apply_italic_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.apply_align_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row2, text="对齐", variable=self.apply_align_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        tk.Button(batch_frame, text="应用到选中框", command=self.apply_style_to_selected,
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 9),
                 cursor="hand2", relief=tk.FLAT).pack(fill=tk.X, pady=5)

        # === 对齐工具 ===
        self.create_section_header(self.prop_frame, "多框对齐")

        align_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        align_frame.pack(fill=tk.X, padx=10, pady=5)

        # 全选按钮
        select_all_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        select_all_frame.pack(fill=tk.X, pady=(0, 5))

        tk.Button(select_all_frame, text="全选当前页所有框 (Ctrl+A)", command=self.select_all_boxes,
                 bg="#FF9800", fg="white", font=(FONT_FAMILY, 9, "bold"),
                 cursor="hand2", relief=tk.FLAT).pack(fill=tk.X)

        tk.Label(align_frame, text="Ctrl+点击选中多个框：",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w", pady=(5, 0))

        # 水平对齐
        h_align_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        h_align_frame.pack(fill=tk.X, pady=3)

        tk.Label(h_align_frame, text="水平:", bg=COLOR_WHITE, font=(FONT_FAMILY, 8), fg="#666").pack(side=tk.LEFT)

        tk.Button(h_align_frame, text="左", command=lambda: self.align_boxes("left"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(h_align_frame, text="中", command=lambda: self.align_boxes("center_h"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(h_align_frame, text="右", command=lambda: self.align_boxes("right"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 垂直对齐
        v_align_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        v_align_frame.pack(fill=tk.X, pady=3)

        tk.Label(v_align_frame, text="垂直:", bg=COLOR_WHITE, font=(FONT_FAMILY, 8), fg="#666").pack(side=tk.LEFT)

        tk.Button(v_align_frame, text="上", command=lambda: self.align_boxes("top"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(v_align_frame, text="中", command=lambda: self.align_boxes("center_v"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(v_align_frame, text="下", command=lambda: self.align_boxes("bottom"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 分隔线
        tk.Frame(align_frame, bg="#e0e0e0", height=1).pack(fill=tk.X, pady=8)

        # 均匀分布
        tk.Label(align_frame, text="均匀分布（需要3个或以上）：",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")

        dist_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        dist_frame.pack(fill=tk.X, pady=3)

        tk.Button(dist_frame, text="水平等间距", command=lambda: self.distribute_boxes("horizontal"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 8), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(dist_frame, text="垂直等间距", command=lambda: self.distribute_boxes("vertical"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 8), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 分隔线
        tk.Frame(align_frame, bg="#e0e0e0", height=1).pack(fill=tk.X, pady=8)

        # 尺寸统一
        tk.Label(align_frame, text="尺寸统一（以第一个选中框为基准）：",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")

        size_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        size_frame.pack(fill=tk.X, pady=3)

        tk.Button(size_frame, text="统一宽", command=lambda: self.unify_size("width"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 8), width=7,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(size_frame, text="统一高", command=lambda: self.unify_size("height"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 8), width=7,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(size_frame, text="统一大小", command=lambda: self.unify_size("both"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 8), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 分隔线
        tk.Frame(align_frame, bg="#e0e0e0", height=1).pack(fill=tk.X, pady=8)

        # 对齐到画布
        tk.Label(align_frame, text="对齐到画布中心：",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")

        canvas_align_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        canvas_align_frame.pack(fill=tk.X, pady=3)

        tk.Button(canvas_align_frame, text="水平居中", command=lambda: self.align_to_canvas("h"),
                 bg="#D32F2F", fg="white", font=(FONT_FAMILY, 8), width=9,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(canvas_align_frame, text="垂直居中", command=lambda: self.align_to_canvas("v"),
                 bg="#D32F2F", fg="white", font=(FONT_FAMILY, 8), width=9,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(canvas_align_frame, text="完全居中", command=lambda: self.align_to_canvas("center"),
                 bg="#D32F2F", fg="white", font=(FONT_FAMILY, 8), width=9,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 批量位移
        tk.Frame(align_frame, bg="#e0e0e0", height=1).pack(fill=tk.X, pady=8)

        tk.Label(align_frame, text="批量位移（像素）：",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")

        # 位移输入框
        offset_input_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        offset_input_frame.pack(fill=tk.X, pady=3)

        tk.Label(offset_input_frame, text="移动:", bg=COLOR_WHITE, font=(FONT_FAMILY, 8), fg="#666").pack(side=tk.LEFT)

        self.offset_px_var = tk.StringVar(value="10")
        offset_entry = tk.Entry(offset_input_frame, textvariable=self.offset_px_var,
                               width=5, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        offset_entry.pack(side=tk.LEFT, padx=3)

        tk.Label(offset_input_frame, text="px", bg=COLOR_WHITE, font=(FONT_FAMILY, 8), fg="#666").pack(side=tk.LEFT)

        # 方向按钮
        offset_btn_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        offset_btn_frame.pack(fill=tk.X, pady=3)

        # 上按钮
        tk.Button(offset_btn_frame, text="↑", command=lambda: self.batch_offset(0, -1),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"), width=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=1)

        # 下按钮
        tk.Button(offset_btn_frame, text="↓", command=lambda: self.batch_offset(0, 1),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"), width=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=1)

        # 左按钮
        tk.Button(offset_btn_frame, text="←", command=lambda: self.batch_offset(-1, 0),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"), width=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=1)

        # 右按钮
        tk.Button(offset_btn_frame, text="→", command=lambda: self.batch_offset(1, 0),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"), width=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=1)

        # === 当前页背景 ===
        self.create_section_header(self.prop_frame, "当前页背景")

        bg_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        bg_frame.pack(fill=tk.X, padx=10, pady=5)

        tk.Label(bg_frame, text="背景图会自动调整为与原图相同大小",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8), wraplength=220).pack(anchor="w")

        bg_btn_frame = tk.Frame(bg_frame, bg=COLOR_WHITE)
        bg_btn_frame.pack(fill=tk.X, pady=5)

        tk.Button(bg_btn_frame, text="设置背景", command=self.load_current_page_background,
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9),
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        tk.Button(bg_btn_frame, text="清除背景", command=self.clear_current_page_background,
                 bg=COLOR_RED, fg="white", font=(FONT_FAMILY, 9),
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 背景状态显示
        self.bg_status_label = tk.Label(bg_frame, text="未设置背景",
                                        bg=COLOR_WHITE, fg="#999", font=(FONT_FAMILY, 8))
        self.bg_status_label.pack(anchor="w", pady=2)

    def create_section_header(self, parent, text):
        """创建属性面板分组标题"""
        header = tk.Frame(parent, bg="#e3f2fd")
        header.pack(fill=tk.X, pady=(10, 5))

        label = tk.Label(header, text=text, bg="#e3f2fd", fg="#1565C0",
                        font=(FONT_FAMILY, 9, "bold"), padx=10, pady=3)
        label.pack(fill=tk.X)

    def create_status_bar(self):
        """创建底部状态栏 - PowerPoint红色主题"""
        self.status_bar = tk.Frame(self.root, bg=COLOR_THEME, height=28)
        self.status_bar.pack(fill=tk.X, side=tk.BOTTOM)
        self.status_bar.pack_propagate(False)

        self.status_label = tk.Label(self.status_bar, text="就绪 - 请导入图片开始编辑",
                                     bg=COLOR_THEME, fg="white",
                                     font=(FONT_FAMILY, 9), padx=10)
        self.status_label.pack(side=tk.LEFT)

        self.status_info = tk.Label(self.status_bar, text="",
                                    bg=COLOR_THEME, fg="white",
                                    font=(FONT_FAMILY, 9), padx=10)
        self.status_info.pack(side=tk.RIGHT)

    def update_status(self, text):
        """更新状态栏"""
        self.status_label.config(text=text)

    def bind_shortcuts(self):
        """绑定快捷键"""
        self.root.bind("<Control-z>", lambda e: self.undo())
        self.root.bind("<Control-y>", lambda e: self.redo())
        self.root.bind("<Delete>", lambda e: self.delete_selected_box())
        self.root.bind("<Left>", lambda e: self.prev_page())
        self.root.bind("<Right>", lambda e: self.next_page())
        self.root.bind("<Control-s>", lambda e: self.save_project())
        self.root.bind("<Control-o>", lambda e: self.load_project())
        # 新增快捷键
        self.root.bind("<Control-a>", lambda e: self.select_all_boxes())
        self.root.bind("<Control-c>", lambda e: self.copy_boxes())
        self.root.bind("<Control-v>", lambda e: self.paste_boxes())
        self.root.bind("<Left>", lambda e: self.move_box_by_key(-10, 0))
        self.root.bind("<Right>", lambda e: self.move_box_by_key(10, 0))
        self.root.bind("<Up>", lambda e: self.move_box_by_key(0, -10))
        self.root.bind("<Down>", lambda e: self.move_box_by_key(0, 10))
        self.root.bind("<Control-Left>", lambda e: self.move_box_by_key(-1, 0))
        self.root.bind("<Control-Right>", lambda e: self.move_box_by_key(1, 0))
        self.root.bind("<Control-Up>", lambda e: self.move_box_by_key(0, -1))
        self.root.bind("<Control-Down>", lambda e: self.move_box_by_key(0, 1))
        self.root.bind("<Prior>", lambda e: self.prev_page())
        self.root.bind("<Next>", lambda e: self.next_page())


    # ==================== 页面管理 ====================

    # 编辑用的最大图片尺寸（超过此尺寸会缩放以提高性能）
    MAX_EDIT_SIZE = 1920

    def _resize_image_for_edit(self, img):
        """缩放图片用于编辑，返回缩放后的图片和缩放比例"""
        w, h = img.size
        if max(w, h) <= self.MAX_EDIT_SIZE:
            return img, 1.0

        scale = self.MAX_EDIT_SIZE / max(w, h)
        new_w = int(w * scale)
        new_h = int(h * scale)
        resized = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
        return resized, scale

    def load_multiple_images(self):
        """批量加载多张原图"""
        file_paths = filedialog.askopenfilenames(
            title="选择多张原图（按顺序选择）",
            filetypes=[("图片文件", "*.jpg *.jpeg *.png *.bmp")]
        )
        if not file_paths:
            return

        if self.pages:
            self.save_current_page()

        clear_existing = False
        if self.pages:
            result = messagebox.askyesnocancel(
                "提示", f"已有 {len(self.pages)} 页，是否清空？\n\n是 - 清空后导入\n否 - 追加\n取消 - 取消"
            )
            if result is None:
                return
            elif result:
                self.pages = []
                clear_existing = True

        start_index = len(self.pages)

        for path in file_paths:
            original_img = Image.open(path)
            original_size = original_img.size  # 保存原始尺寸

            # 缩放图片用于编辑
            edit_img, edit_scale = self._resize_image_for_edit(original_img)

            page_data = {
                "original_path": path,
                "original_size": original_size,  # 原始尺寸
                "edit_scale": edit_scale,  # 编辑缩放比例
                "bg_path": None,
                "bg_original_path": None,  # 背景原图路径
                "image": edit_img,  # 编辑用的缩放图片
                "text_boxes": []
            }
            self.pages.append(page_data)

        self.current_page_index = start_index
        self.load_current_page()
        self.update_page_label()
        self.update_thumbnails()

        # 隐藏占位符
        self.placeholder_label.place_forget()

        # 显示是否有缩放
        any_scaled = any(p["edit_scale"] < 1.0 for p in self.pages[start_index:])
        if any_scaled:
            self.update_status(f"已导入 {len(file_paths)} 张图片（大图已自动缩放以提高性能），共 {len(self.pages)} 页")
        else:
            self.update_status(f"已导入 {len(file_paths)} 张图片，共 {len(self.pages)} 页")

    def load_multiple_backgrounds(self):
        """批量加载背景图 - 自动调整大小与编辑图一致"""
        if not self.pages:
            messagebox.showwarning("提示", "请先导入原图")
            return

        file_paths = filedialog.askopenfilenames(
            title="选择背景图",
            filetypes=[("图片文件", "*.jpg *.jpeg *.png *.bmp")]
        )
        if not file_paths:
            return

        matched = 0
        for bg_path in file_paths:
            bg_name = os.path.splitext(os.path.basename(bg_path))[0].lower()
            for page in self.pages:
                orig_name = os.path.splitext(os.path.basename(page["original_path"]))[0].lower()
                if bg_name == orig_name or orig_name in bg_name or bg_name in orig_name:
                    # 调整背景图大小与编辑图一致
                    resized_bg_path = self._resize_bg_to_match(bg_path, page["image"].size)
                    page["bg_path"] = resized_bg_path
                    matched += 1
                    break

        # 如果没有匹配到，按顺序分配
        if matched == 0 and len(file_paths) == len(self.pages):
            for i, bg_path in enumerate(file_paths):
                resized_bg_path = self._resize_bg_to_match(bg_path, self.pages[i]["image"].size)
                self.pages[i]["bg_path"] = resized_bg_path
            matched = len(file_paths)

        # 更新当前页背景路径
        if self.pages and self.current_page_index < len(self.pages):
            self.clean_bg_path = self.pages[self.current_page_index].get("bg_path")

        # 刷新显示
        self.update_bg_status()
        self.update_thumbnails()
        self.refresh_canvas()
        self.update_status(f"已匹配 {matched}/{len(self.pages)} 张背景图")

    def load_current_page_background(self):
        """为当前页单独设置背景图"""
        if not self.pages:
            messagebox.showwarning("提示", "请先导入原图")
            return

        file_path = filedialog.askopenfilename(
            title=f"选择第 {self.current_page_index + 1} 页的背景图",
            filetypes=[("图片文件", "*.jpg *.jpeg *.png *.bmp")]
        )
        if not file_path:
            return

        page = self.pages[self.current_page_index]
        edit_size = page["image"].size

        # 调整背景图大小与编辑图一致
        resized_bg_path = self._resize_bg_to_match(file_path, edit_size)
        page["bg_path"] = resized_bg_path
        self.clean_bg_path = resized_bg_path

        self.update_bg_status()
        self.update_thumbnails()
        self.refresh_canvas()
        self.update_status(f"第 {self.current_page_index + 1} 页背景已设置")

    def _resize_bg_to_match(self, bg_path, target_size):
        """调整背景图大小与目标尺寸一致，返回调整后的图片路径"""
        bg_img = Image.open(bg_path)

        # 如果大小已经一致，直接返回原路径
        if bg_img.size == target_size:
            return bg_path

        # 调整大小
        resized_img = bg_img.resize(target_size, Image.Resampling.LANCZOS)

        # 保存到临时文件
        bg_dir = os.path.dirname(bg_path)
        bg_name = os.path.splitext(os.path.basename(bg_path))[0]
        bg_ext = os.path.splitext(bg_path)[1]

        # 创建调整后的文件名
        resized_path = os.path.join(bg_dir, f"{bg_name}_resized_{target_size[0]}x{target_size[1]}{bg_ext}")

        # 如果已存在同名调整后的文件，检查是否需要重新生成
        if not os.path.exists(resized_path):
            if resized_img.mode == 'RGBA' and bg_ext.lower() in ['.jpg', '.jpeg']:
                resized_img = resized_img.convert('RGB')
            resized_img.save(resized_path, quality=95)

        return resized_path

    def clear_current_page_background(self):
        """清除当前页背景"""
        if not self.pages:
            return

        self.pages[self.current_page_index]["bg_path"] = None
        self.clean_bg_path = None
        self.update_bg_status()
        self.update_thumbnails()
        self.refresh_canvas()
        self.update_status(f"第 {self.current_page_index + 1} 页背景已清除")

    def update_bg_status(self):
        """更新背景状态显示"""
        if self.clean_bg_path:
            bg_name = os.path.basename(self.clean_bg_path)
            if len(bg_name) > 25:
                bg_name = bg_name[:22] + "..."
            self.bg_status_label.config(text=f"已设置: {bg_name}", fg=COLOR_GREEN)
        else:
            self.bg_status_label.config(text="未设置背景", fg="#999")

    def save_current_page(self):
        """保存当前页数据"""
        if not self.pages or self.current_page_index >= len(self.pages):
            return
        page = self.pages[self.current_page_index]
        page["text_boxes"] = [box.to_dict() for box in self.text_boxes]
        page["bg_path"] = self.clean_bg_path

    def load_current_page(self):
        """加载当前页数据"""
        if not self.pages or self.current_page_index >= len(self.pages):
            return

        page = self.pages[self.current_page_index]
        self.original_img_path = page["original_path"]
        self.original_image = page["image"]
        self.clean_bg_path = page.get("bg_path")
        self.text_boxes = [TextBox.from_dict(d) for d in page.get("text_boxes", [])]
        self.selected_box_index = -1
        self.selected_boxes = []

        self.fit_image_to_canvas()
        self.update_listbox()
        self.update_status_info()
        self.update_bg_status()

    def prev_page(self):
        """上一页"""
        if not self.pages or self.current_page_index <= 0:
            return
        self.save_current_page()
        self.current_page_index -= 1
        self.load_current_page()
        self.update_page_label()
        self.highlight_current_thumbnail()

    def next_page(self):
        """下一页"""
        if not self.pages or self.current_page_index >= len(self.pages) - 1:
            return
        self.save_current_page()
        self.current_page_index += 1
        self.load_current_page()
        self.update_page_label()
        self.highlight_current_thumbnail()

    def update_page_label(self):
        """更新页码"""
        if self.pages:
            page_text = f"{self.current_page_index + 1}/{len(self.pages)}"
            self.page_label.config(text=page_text)
            self.title_page_label.config(text=f"第 {page_text} 页")
        else:
            self.page_label.config(text="0/0")
            self.title_page_label.config(text="第 0/0 页")

    def update_status_info(self):
        """更新状态栏信息"""
        if self.pages and self.original_image:
            w, h = self.original_image.size
            boxes = len(self.text_boxes)
            self.status_info.config(text=f"尺寸: {w}×{h} | 文本框: {boxes} | 缩放: {int(self.scale*100)}%")
            self.zoom_label.config(text=f"{int(self.scale*100)}%")

    def update_thumbnails(self):
        """更新缩略图"""
        # 清空现有缩略图
        for widget in self.thumbnail_frame.winfo_children():
            widget.destroy()
        self.thumbnail_images = []

        for idx, page in enumerate(self.pages):
            frame = tk.Frame(self.thumbnail_frame, bg="#ffffff", cursor="hand2",
                           relief=tk.GROOVE, bd=1)
            frame.pack(fill=tk.X, padx=5, pady=3)

            # 生成缩略图
            img = page["image"].copy()
            img.thumbnail((110, 70), Image.Resampling.LANCZOS)
            tk_img = ImageTk.PhotoImage(img)
            self.thumbnail_images.append(tk_img)

            # 缩略图标签
            label = tk.Label(frame, image=tk_img, bg="#ffffff")
            label.pack(padx=2, pady=2)

            # 页码和背景状态
            has_bg = "✓" if page.get("bg_path") else ""
            page_num = tk.Label(frame, text=f"第 {idx + 1} 页 {has_bg}", bg="#ffffff",
                               fg="#666666" if not has_bg else COLOR_GREEN, font=("微软雅黑", 8))
            page_num.pack()

            # 点击切换页面
            frame.bind("<Button-1>", lambda e, i=idx: self.go_to_page(i))
            label.bind("<Button-1>", lambda e, i=idx: self.go_to_page(i))
            page_num.bind("<Button-1>", lambda e, i=idx: self.go_to_page(i))

            # 右键菜单
            frame.bind("<Button-3>", lambda e, i=idx: self.show_thumbnail_menu(e, i))
            label.bind("<Button-3>", lambda e, i=idx: self.show_thumbnail_menu(e, i))
            page_num.bind("<Button-3>", lambda e, i=idx: self.show_thumbnail_menu(e, i))

        self.highlight_current_thumbnail()

    def show_thumbnail_menu(self, event, page_index):
        """显示缩略图右键菜单"""
        menu = tk.Menu(self.root, tearoff=0)
        menu.add_command(label=f"设置第 {page_index + 1} 页背景",
                        command=lambda: self.set_page_background(page_index))
        menu.add_command(label=f"清除第 {page_index + 1} 页背景",
                        command=lambda: self.clear_page_background(page_index))
        menu.add_separator()
        menu.add_command(label=f"删除第 {page_index + 1} 页",
                        command=lambda: self.delete_page(page_index))
        menu.post(event.x_root, event.y_root)

    def set_page_background(self, page_index):
        """为指定页设置背景图"""
        if page_index < 0 or page_index >= len(self.pages):
            return

        file_path = filedialog.askopenfilename(
            title=f"选择第 {page_index + 1} 页的背景图",
            filetypes=[("图片文件", "*.jpg *.jpeg *.png *.bmp")]
        )
        if not file_path:
            return

        page = self.pages[page_index]
        edit_size = page["image"].size

        # 调整背景图大小与编辑图一致
        resized_bg_path = self._resize_bg_to_match(file_path, edit_size)
        page["bg_path"] = resized_bg_path

        # 如果是当前页，更新当前页的背景路径并刷新画布
        if page_index == self.current_page_index:
            self.clean_bg_path = resized_bg_path
            self.update_bg_status()
            self.refresh_canvas()

        # 更新缩略图显示
        self.update_thumbnails()
        self.update_status(f"第 {page_index + 1} 页背景已设置")

    def clear_page_background(self, page_index):
        """清除指定页的背景图"""
        if page_index < 0 or page_index >= len(self.pages):
            return

        self.pages[page_index]["bg_path"] = None

        # 如果是当前页，更新当前页的背景路径并刷新画布
        if page_index == self.current_page_index:
            self.clean_bg_path = None
            self.update_bg_status()
            self.refresh_canvas()

        # 更新缩略图显示
        self.update_thumbnails()
        self.update_status(f"第 {page_index + 1} 页背景已清除")

    def delete_page(self, page_index):
        """删除指定页"""
        if page_index < 0 or page_index >= len(self.pages):
            return

        if len(self.pages) <= 1:
            messagebox.showwarning("提示", "至少保留一页")
            return

        result = messagebox.askyesno("确认", f"确定删除第 {page_index + 1} 页？")
        if not result:
            return

        del self.pages[page_index]

        # 调整当前页索引
        if self.current_page_index >= len(self.pages):
            self.current_page_index = len(self.pages) - 1
        elif self.current_page_index > page_index:
            self.current_page_index -= 1

        self.load_current_page()
        self.update_page_label()
        self.update_thumbnails()
        self.update_status(f"已删除页面，剩余 {len(self.pages)} 页")

    def highlight_current_thumbnail(self):
        """高亮当前页缩略图"""
        for idx, widget in enumerate(self.thumbnail_frame.winfo_children()):
            if idx == self.current_page_index:
                widget.config(bg="#bbdefb", relief=tk.SOLID, bd=2)
                for child in widget.winfo_children():
                    child.config(bg="#bbdefb")
            else:
                widget.config(bg="#ffffff", relief=tk.GROOVE, bd=1)
                for child in widget.winfo_children():
                    child.config(bg="#ffffff")

    def go_to_page(self, index):
        """跳转到指定页"""
        if 0 <= index < len(self.pages):
            self.save_current_page()
            self.current_page_index = index
            self.load_current_page()
            self.update_page_label()
            self.highlight_current_thumbnail()

    # ==================== 画布操作 ====================

    def fit_image_to_canvas(self):
        """自适应显示图片"""
        if not self.original_image:
            return

        canvas_w = self.canvas.winfo_width()
        canvas_h = self.canvas.winfo_height()

        if canvas_w < 10 or canvas_h < 10:
            self.root.after(100, self.fit_image_to_canvas)
            return

        img_w, img_h = self.original_image.size
        scale_w = (canvas_w - 40) / img_w
        scale_h = (canvas_h - 40) / img_h
        self.scale = min(scale_w, scale_h, 1.0)

        self.refresh_canvas()

    def on_canvas_resize(self, event):
        """画布大小改变"""
        if self.original_image:
            self.fit_image_to_canvas()

    def on_canvas_zoom(self, event):
        """Ctrl+滚轮缩放"""
        if not self.original_image:
            return

        # 获取鼠标位置作为缩放中心
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        # 缩放因子
        if event.delta > 0:
            factor = 1.1  # 放大
        else:
            factor = 0.9  # 缩小

        # 计算新缩放比例（限制范围 10% - 300%）
        new_scale = self.scale * factor
        new_scale = max(0.1, min(new_scale, 3.0))

        if new_scale != self.scale:
            self.scale = new_scale
            self.refresh_canvas()
            self.update_status(f"缩放: {int(self.scale * 100)}%")

    def on_canvas_scroll(self, event):
        """普通滚轮滚动"""
        self.canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")

    def zoom_to_100(self):
        """缩放到100%"""
        if not self.original_image:
            return
        self.scale = 1.0
        self.refresh_canvas()
        self.update_status("缩放: 100%")

    def refresh_canvas(self):
        """刷新画布"""
        if not self.original_image:
            return

        if self.current_preview_mode == "ppt":
            self._draw_ppt_preview()
        else:
            self._draw_original_with_boxes()

        self.update_status_info()

    def _draw_original_with_boxes(self):
        """绘制原图+框"""
        self.canvas.delete("all")

        img_w, img_h = self.original_image.size
        display_w = int(img_w * self.scale)
        display_h = int(img_h * self.scale)

        # 居中显示
        canvas_w = self.canvas.winfo_width()
        canvas_h = self.canvas.winfo_height()
        offset_x = max(0, (canvas_w - display_w) // 2)
        offset_y = max(0, (canvas_h - display_h) // 2)

        self.display_image = self.original_image.resize((display_w, display_h), Image.Resampling.LANCZOS)
        self.tk_image = ImageTk.PhotoImage(self.display_image)
        self.canvas.create_image(offset_x, offset_y, anchor=tk.NW, image=self.tk_image, tags="image")

        # 保存偏移量用于坐标转换
        self.canvas_offset_x = offset_x
        self.canvas_offset_y = offset_y

        for idx, box in enumerate(self.text_boxes):
            self.draw_box(idx, box, offset_x, offset_y)

        self.canvas.config(scrollregion=(0, 0, max(canvas_w, display_w + offset_x * 2),
                                          max(canvas_h, display_h + offset_y * 2)))

    def _draw_ppt_preview(self):
        """绘制PPT预览"""
        if self.clean_bg_path and os.path.exists(self.clean_bg_path):
            bg_image = Image.open(self.clean_bg_path)
        elif self.original_image:
            bg_image = self.original_image.copy()
        else:
            return

        preview_img = bg_image.copy().convert("RGBA")
        img_w, img_h = preview_img.size

        try:
            draw = ImageDraw.Draw(preview_img)

            for box in self.text_boxes:
                if not box.text:
                    continue

                pixel_font_size = int(box.font_size * 96 / 72)

                try:
                    font_path = self._get_font_path(box.font_name)
                    if font_path:
                        font = ImageFont.truetype(font_path, pixel_font_size)
                    else:
                        font = ImageFont.load_default()
                except:
                    font = ImageFont.load_default()

                color_hex = box.font_color.lstrip('#')
                r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)

                text_x, text_y = box.x + 3, box.y + 2

                try:
                    bbox = draw.textbbox((0, 0), box.text, font=font)
                    text_width = bbox[2] - bbox[0]
                    text_height = bbox[3] - bbox[1]
                except:
                    text_width = len(box.text) * pixel_font_size
                    text_height = pixel_font_size

                if box.align == "center":
                    text_x = box.x + (box.width - text_width) // 2
                elif box.align == "right":
                    text_x = box.x + box.width - text_width - 3

                text_y = box.y + (box.height - text_height) // 2

                draw.text((text_x, text_y), box.text, font=font, fill=(r, g, b, 255))

        except Exception as e:
            print(f"绘制文字失败: {e}")

        preview_img = preview_img.convert("RGB")

        canvas_w = self.canvas.winfo_width()
        canvas_h = self.canvas.winfo_height()

        # 保持当前缩放比例，不强制重置
        display_w = int(img_w * self.scale)
        display_h = int(img_h * self.scale)

        offset_x = max(0, (canvas_w - display_w) // 2)
        offset_y = max(0, (canvas_h - display_h) // 2)

        self.canvas_offset_x = offset_x
        self.canvas_offset_y = offset_y

        preview_img = preview_img.resize((display_w, display_h), Image.Resampling.LANCZOS)
        self.ppt_preview_image = ImageTk.PhotoImage(preview_img)

        self.canvas.delete("all")
        self.canvas.create_image(offset_x, offset_y, anchor=tk.NW, image=self.ppt_preview_image)

        for idx, box in enumerate(self.text_boxes):
            self._draw_ppt_edit_box(idx, box, offset_x, offset_y)

        self.canvas.config(scrollregion=(0, 0, max(canvas_w, display_w + offset_x * 2),
                                          max(canvas_h, display_h + offset_y * 2)))

    def _get_font_path(self, font_name):
        """获取字体路径"""
        font_map = {
            "微软雅黑": "C:/Windows/Fonts/msyh.ttc",
            "宋体": "C:/Windows/Fonts/simsun.ttc",
            "黑体": "C:/Windows/Fonts/simhei.ttf",
            "楷体": "C:/Windows/Fonts/simkai.ttf",
            "仿宋": "C:/Windows/Fonts/simfang.ttf",
            "Arial": "C:/Windows/Fonts/arial.ttf"
        }
        path = font_map.get(font_name)
        if path and os.path.exists(path):
            return path
        return font_map.get("微软雅黑")

    def draw_box(self, idx, box, offset_x, offset_y):
        """绘制文本框"""
        x1 = int(box.x * self.scale) + offset_x
        y1 = int(box.y * self.scale) + offset_y
        x2 = int((box.x + box.width) * self.scale) + offset_x
        y2 = int((box.y + box.height) * self.scale) + offset_y

        is_primary = (idx == self.selected_box_index)
        is_multi = (idx in self.selected_boxes)

        if is_primary:
            color, width = "#1976D2", 3
        elif is_multi:
            color, width = "#4CAF50", 2
        else:
            color, width = "#f44336", 2

        self.canvas.create_rectangle(x1, y1, x2, y2, outline=color, width=width, tags=f"box_{idx}")

        # 序号
        self.canvas.create_oval(x1 + 5, y1 + 5, x1 + 22, y1 + 22, fill="#FF9800", outline="")
        self.canvas.create_text(x1 + 13, y1 + 13, text=str(idx + 1), fill="white", font=("Arial", 8, "bold"))

        # 文本预览
        if box.text and y2 - y1 > 30:
            preview = box.text[:15] + "..." if len(box.text) > 15 else box.text
            self.canvas.create_text(x1 + 5, y2 - 12, text=preview, fill="#333333",
                                   anchor=tk.NW, font=("微软雅黑", 8))

        # 选中手柄
        if is_primary:
            handle_size = 8
            handles = [(x1, y1), (x2, y1), (x1, y2), (x2, y2),
                      ((x1+x2)//2, y1), ((x1+x2)//2, y2), (x1, (y1+y2)//2), (x2, (y1+y2)//2)]
            for hx, hy in handles:
                self.canvas.create_rectangle(hx - handle_size//2, hy - handle_size//2,
                                            hx + handle_size//2, hy + handle_size//2,
                                            fill="#1976D2", outline="white")

    def _draw_ppt_edit_box(self, idx, box, offset_x, offset_y):
        """PPT预览模式下的编辑框"""
        x1 = int(box.x * self.scale) + offset_x
        y1 = int(box.y * self.scale) + offset_y
        x2 = int((box.x + box.width) * self.scale) + offset_x
        y2 = int((box.y + box.height) * self.scale) + offset_y

        is_primary = (idx == self.selected_box_index)
        is_multi = (idx in self.selected_boxes)

        if is_primary:
            self.canvas.create_rectangle(x1, y1, x2, y2, outline="#1976D2", width=2, dash=(4, 4))
            handle_size = 8
            handles = [(x1, y1), (x2, y1), (x1, y2), (x2, y2),
                      ((x1+x2)//2, y1), ((x1+x2)//2, y2), (x1, (y1+y2)//2), (x2, (y1+y2)//2)]
            for hx, hy in handles:
                self.canvas.create_rectangle(hx - handle_size//2, hy - handle_size//2,
                                            hx + handle_size//2, hy + handle_size//2,
                                            fill="#1976D2", outline="white")
        elif is_multi:
            self.canvas.create_rectangle(x1, y1, x2, y2, outline="#4CAF50", width=2, dash=(4, 4))
        else:
            self.canvas.create_rectangle(x1, y1, x2, y2, outline="#999999", width=1, dash=(2, 4))

    # ==================== 鼠标事件 ====================

    def on_canvas_press(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        # 转换为图片坐标
        img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
        img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale

        if self.selected_box_index >= 0:
            handle = self.check_resize_handle(canvas_x, canvas_y)
            if handle:
                self.is_resizing = True
                self.resize_handle = handle
                self.drag_start_x = canvas_x
                self.drag_start_y = canvas_y
                return

        clicked_idx = self.find_box_at(img_x, img_y)

        if clicked_idx >= 0:
            self.select_box(clicked_idx)
            self.is_dragging = True
            self.drag_start_x = canvas_x
            self.drag_start_y = canvas_y
        elif self.draw_mode:
            # 画框模式
            self.is_drawing = True
            self.draw_start_x = img_x
            self.draw_start_y = img_y
        else:
            # 选择模式：开始框选
            self.is_selecting = True
            self.select_start_x = canvas_x
            self.select_start_y = canvas_y

    def on_canvas_ctrl_click(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
        img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale

        clicked_idx = self.find_box_at(img_x, img_y)

        if clicked_idx >= 0:
            if clicked_idx in self.selected_boxes:
                self.selected_boxes.remove(clicked_idx)
            else:
                self.selected_boxes.append(clicked_idx)

            if self.selected_boxes:
                self.selected_box_index = self.selected_boxes[-1]
            else:
                self.selected_box_index = -1

            self.refresh_canvas()
            self.update_property_panel()
            self.update_status(f"已选中 {len(self.selected_boxes)} 个框")

    def on_canvas_drag(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        if self.is_resizing and self.selected_box_index >= 0:
            self.resize_selected_box(canvas_x, canvas_y)
        elif self.is_dragging and self.selected_box_index >= 0:
            self.drag_selected_box(canvas_x, canvas_y)
        elif self.is_drawing:
            self.draw_temp_rect(canvas_x, canvas_y)
        elif self.is_selecting:
            self.draw_selection_rect(canvas_x, canvas_y)

    def on_canvas_release(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        if self.is_drawing:
            self.finish_drawing(canvas_x, canvas_y)
        elif self.is_selecting:
            self.finish_selection(canvas_x, canvas_y)

        self.is_drawing = False
        self.is_dragging = False
        self.is_resizing = False
        self.is_selecting = False
        self.resize_handle = None
        self.canvas.delete("temp_rect")
        self.canvas.delete("selection_rect")

    def on_canvas_double_click(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
        img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale

        clicked_idx = self.find_box_at(img_x, img_y)
        if clicked_idx >= 0:
            self.select_box(clicked_idx)
            self.show_inline_text_editor(clicked_idx)

    def on_canvas_right_click(self, event):
        """右键菜单"""
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
        img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale

        # 查找点击的文本框
        clicked_idx = self.find_box_at(img_x, img_y)

        # 创建右键菜单
        menu = tk.Menu(self.root, tearoff=0, font=(FONT_FAMILY, 9))

        if clicked_idx >= 0:
            # 点击在文本框上
            self.select_box(clicked_idx)

            menu.add_command(label="🔍 OCR识别此框", command=self.ocr_single_box,
                           font=(FONT_FAMILY, 9, "bold"))
            menu.add_separator()
            menu.add_command(label="✏️ 编辑文字", command=lambda: self.show_inline_text_editor(clicked_idx))
            menu.add_separator()
            menu.add_command(label="📋 复制 (Ctrl+C)", command=self.copy_boxes)
            menu.add_command(label="📄 粘贴 (Ctrl+V)", command=self.paste_boxes)
            menu.add_separator()
            menu.add_command(label="🗑️ 删除 (Del)", command=self.delete_selected_box,
                           foreground=COLOR_RED)
        else:
            # 点击在空白处
            if self.clipboard_boxes:
                menu.add_command(label="📄 粘贴 (Ctrl+V)", command=self.paste_boxes)
                menu.add_separator()

            menu.add_command(label="📐 开始画框", command=self.toggle_draw_mode_btn)

            if self.text_boxes:
                menu.add_separator()
                menu.add_command(label="🔍 OCR识别全部", command=self.ocr_all_boxes)

        # 显示菜单
        try:
            menu.tk_popup(event.x_root, event.y_root)
        finally:
            menu.grab_release()

    def find_box_at(self, x, y):
        for idx in range(len(self.text_boxes) - 1, -1, -1):
            box = self.text_boxes[idx]
            if box.x <= x <= box.x + box.width and box.y <= y <= box.y + box.height:
                return idx
        return -1

    def check_resize_handle(self, canvas_x, canvas_y):
        if self.selected_box_index < 0:
            return None

        box = self.text_boxes[self.selected_box_index]
        offset_x = getattr(self, 'canvas_offset_x', 0)
        offset_y = getattr(self, 'canvas_offset_y', 0)

        x1 = int(box.x * self.scale) + offset_x
        y1 = int(box.y * self.scale) + offset_y
        x2 = int((box.x + box.width) * self.scale) + offset_x
        y2 = int((box.y + box.height) * self.scale) + offset_y

        handle_size = 10
        handles = {
            "nw": (x1, y1), "ne": (x2, y1), "sw": (x1, y2), "se": (x2, y2),
            "n": ((x1+x2)//2, y1), "s": ((x1+x2)//2, y2),
            "w": (x1, (y1+y2)//2), "e": (x2, (y1+y2)//2)
        }

        for handle_type, (hx, hy) in handles.items():
            if abs(canvas_x - hx) < handle_size and abs(canvas_y - hy) < handle_size:
                return handle_type
        return None

    def draw_temp_rect(self, canvas_x, canvas_y):
        self.canvas.delete("temp_rect")

        offset_x = getattr(self, 'canvas_offset_x', 0)
        offset_y = getattr(self, 'canvas_offset_y', 0)

        x1 = int(self.draw_start_x * self.scale) + offset_x
        y1 = int(self.draw_start_y * self.scale) + offset_y
        x2 = int(canvas_x)
        y2 = int(canvas_y)

        self.canvas.create_rectangle(x1, y1, x2, y2, outline="#1976D2", width=2,
                                    dash=(5, 5), tags="temp_rect")

    def finish_drawing(self, canvas_x, canvas_y):
        offset_x = getattr(self, 'canvas_offset_x', 0)
        offset_y = getattr(self, 'canvas_offset_y', 0)

        x1 = self.draw_start_x
        y1 = self.draw_start_y
        x2 = (canvas_x - offset_x) / self.scale
        y2 = (canvas_y - offset_y) / self.scale

        if x1 > x2: x1, x2 = x2, x1
        if y1 > y2: y1, y2 = y2, y1

        width = x2 - x1
        height = y2 - y1

        if width < 10 or height < 10:
            return

        self.save_state()

        new_box = TextBox(int(x1), int(y1), int(width), int(height))
        self.text_boxes.append(new_box)
        self.select_box(len(self.text_boxes) - 1)
        self.refresh_canvas()
        self.update_listbox()
        self.mark_unsaved()

    def draw_selection_rect(self, canvas_x, canvas_y):
        """绘制框选区域"""
        self.canvas.delete("selection_rect")

        x1 = int(self.select_start_x)
        y1 = int(self.select_start_y)
        x2 = int(canvas_x)
        y2 = int(canvas_y)

        # 绘制半透明的蓝色选区矩形
        self.canvas.create_rectangle(x1, y1, x2, y2,
                                     outline="#2196F3", width=2,
                                     dash=(3, 3), tags="selection_rect")

    def finish_selection(self, canvas_x, canvas_y):
        """完成框选，选中选区内的所有框"""
        offset_x = getattr(self, 'canvas_offset_x', 0)
        offset_y = getattr(self, 'canvas_offset_y', 0)

        # 计算选区的图片坐标
        x1 = (self.select_start_x - offset_x) / self.scale
        y1 = (self.select_start_y - offset_y) / self.scale
        x2 = (canvas_x - offset_x) / self.scale
        y2 = (canvas_y - offset_y) / self.scale

        # 确保x1 < x2, y1 < y2
        if x1 > x2: x1, x2 = x2, x1
        if y1 > y2: y1, y2 = y2, y1

        # 选区太小则忽略
        if abs(x2 - x1) < 5 or abs(y2 - y1) < 5:
            return

        # 查找选区内的所有文本框
        selected_indices = []
        for idx, box in enumerate(self.text_boxes):
            # 检查文本框是否与选区相交或包含在选区内
            box_left = box.x
            box_right = box.x + box.width
            box_top = box.y
            box_bottom = box.y + box.height

            # 判断相交：选区的任意部分与框重叠
            if (box_left < x2 and box_right > x1 and
                box_top < y2 and box_bottom > y1):
                selected_indices.append(idx)

        # 选中找到的框
        if selected_indices:
            self.selected_boxes = selected_indices
            self.selected_box_index = selected_indices[0] if selected_indices else -1

            # 更新界面
            self.refresh_canvas()
            self.update_property_panel()

            # 更新列表框选择
            self.box_listbox.selection_clear(0, tk.END)
            for idx in self.selected_boxes:
                self.box_listbox.selection_set(idx)

            self.update_status(f"框选选中 {len(selected_indices)} 个文本框 ✓")
        else:
            # 没有选中任何框，清空选择
            self.selected_boxes = []
            self.selected_box_index = -1
            self.refresh_canvas()
            self.update_status("框选区域内没有文本框")

    def resize_selected_box(self, canvas_x, canvas_y):
        if self.selected_box_index < 0:
            return

        box = self.text_boxes[self.selected_box_index]
        dx = (canvas_x - self.drag_start_x) / self.scale
        dy = (canvas_y - self.drag_start_y) / self.scale

        if "w" in self.resize_handle:
            new_x = box.x + dx
            new_w = box.width - dx
            if new_w > 10:
                box.x = int(new_x)
                box.width = int(new_w)
        if "e" in self.resize_handle:
            new_w = box.width + dx
            if new_w > 10:
                box.width = int(new_w)
        if "n" in self.resize_handle:
            new_y = box.y + dy
            new_h = box.height - dy
            if new_h > 10:
                box.y = int(new_y)
                box.height = int(new_h)
        if "s" in self.resize_handle:
            new_h = box.height + dy
            if new_h > 10:
                box.height = int(new_h)

        self.drag_start_x = canvas_x
        self.drag_start_y = canvas_y
        self.refresh_canvas()
        self.update_property_panel()

    def drag_selected_box(self, canvas_x, canvas_y):
        if self.selected_box_index < 0:
            return

        box = self.text_boxes[self.selected_box_index]
        dx = (canvas_x - self.drag_start_x) / self.scale
        dy = (canvas_y - self.drag_start_y) / self.scale

        box.x = int(box.x + dx)
        box.y = int(box.y + dy)

        self.drag_start_x = canvas_x
        self.drag_start_y = canvas_y
        self.refresh_canvas()
        self.update_property_panel()

    # ==================== 选择与属性 ====================

    def select_box(self, idx):
        self.selected_box_index = idx
        self.selected_boxes = [idx] if idx >= 0 else []
        self.refresh_canvas()
        self.update_property_panel()

        self.box_listbox.selection_clear(0, tk.END)
        if idx >= 0:
            self.box_listbox.selection_set(idx)
            self.box_listbox.see(idx)

    def update_listbox(self):
        self.box_listbox.delete(0, tk.END)
        for idx, box in enumerate(self.text_boxes):
            text_preview = box.text[:15] + "..." if len(box.text) > 15 else box.text
            if not text_preview:
                text_preview = "(空)"
            self.box_listbox.insert(tk.END, f"{idx+1}. {text_preview}")

    def on_listbox_select(self, event):
        selection = self.box_listbox.curselection()
        if selection:
            self.select_box(selection[0])

    def update_property_panel(self):
        if self.selected_box_index < 0 or self.selected_box_index >= len(self.text_boxes):
            return

        box = self.text_boxes[self.selected_box_index]

        self.text_entry.delete("1.0", tk.END)
        self.text_entry.insert("1.0", box.text)

        self.x_entry.delete(0, tk.END)
        self.x_entry.insert(0, str(box.x))
        self.y_entry.delete(0, tk.END)
        self.y_entry.insert(0, str(box.y))
        self.w_entry.delete(0, tk.END)
        self.w_entry.insert(0, str(box.width))
        self.h_entry.delete(0, tk.END)
        self.h_entry.insert(0, str(box.height))

        self.fontsize_var.set(str(box.font_size))
        self.fontname_var.set(box.font_name)
        self.bold_var.set(box.bold)
        self.italic_var.set(box.italic)
        self.align_var.set(box.align)
        self.color_btn.config(bg=box.font_color)

        self.update_style_buttons()
        self.update_align_buttons()

    def update_style_buttons(self):
        if self.bold_var.get():
            self.bold_btn.config(bg="#1976D2", fg="white")
        else:
            self.bold_btn.config(bg="#e0e0e0", fg="black")

        if self.italic_var.get():
            self.italic_btn.config(bg="#1976D2", fg="white")
        else:
            self.italic_btn.config(bg="#e0e0e0", fg="black")

    def on_text_change(self, event=None):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        box.text = self.text_entry.get("1.0", tk.END).strip()
        self.update_listbox()
        self.refresh_canvas()

    def on_position_change(self, event=None):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        try:
            box.x = int(self.x_entry.get())
            box.y = int(self.y_entry.get())
            box.width = int(self.w_entry.get())
            box.height = int(self.h_entry.get())
            self.refresh_canvas()
        except ValueError:
            pass

    def on_font_change(self, event=None):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        try:
            box.font_size = int(self.fontsize_var.get())
        except:
            pass
        box.font_name = self.fontname_var.get()
        self.refresh_canvas()

    def set_align(self, align):
        """设置对齐方式"""
        self.align_var.set(align)
        self.update_align_buttons()
        self.on_style_change()

    def update_align_buttons(self):
        """更新对齐按钮状态"""
        align = self.align_var.get()
        # 左对齐
        if align == "left":
            self.align_left_btn.config(bg="#1976D2", fg="white")
        else:
            self.align_left_btn.config(bg="#e0e0e0", fg="#333")
        # 居中
        if align == "center":
            self.align_center_btn.config(bg="#1976D2", fg="white")
        else:
            self.align_center_btn.config(bg="#e0e0e0", fg="#333")
        # 右对齐
        if align == "right":
            self.align_right_btn.config(bg="#1976D2", fg="white")
        else:
            self.align_right_btn.config(bg="#e0e0e0", fg="#333")

    def on_style_change(self):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        box.bold = self.bold_var.get()
        box.italic = self.italic_var.get()
        box.align = self.align_var.get()
        self.refresh_canvas()

    def toggle_bold(self):
        self.bold_var.set(not self.bold_var.get())
        self.update_style_buttons()
        self.on_style_change()

    def toggle_italic(self):
        self.italic_var.set(not self.italic_var.get())
        self.update_style_buttons()
        self.on_style_change()

    def choose_color(self):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        color = colorchooser.askcolor(color=box.font_color, title="选择文字颜色")
        if color[1]:
            box.font_color = color[1]
            self.color_btn.config(bg=color[1])
            self.refresh_canvas()

    # ==================== 其他操作 ====================

    def toggle_draw_mode(self):
        self.draw_mode = self.draw_mode_var.get()
        if self.draw_mode:
            self.canvas.config(cursor="crosshair")
        else:
            self.canvas.config(cursor="")

    def switch_preview_mode(self):
        self.current_preview_mode = self.preview_mode_var.get()
        self.refresh_canvas()

    def refresh_ppt_preview(self):
        self.preview_mode_var.set("ppt")
        self.current_preview_mode = "ppt"
        self.refresh_canvas()
        self.update_status("PPT预览已刷新 ✓")

    def show_inline_text_editor(self, box_idx):
        """内联文字编辑器"""
        if box_idx < 0 or box_idx >= len(self.text_boxes):
            return

        box = self.text_boxes[box_idx]

        edit_window = tk.Toplevel(self.root)
        edit_window.title(f"编辑文本框 {box_idx + 1}")
        edit_window.geometry("420x300")
        edit_window.configure(bg="#ffffff")
        edit_window.transient(self.root)
        edit_window.grab_set()

        mouse_x = self.root.winfo_pointerx()
        mouse_y = self.root.winfo_pointery()
        edit_window.geometry(f"+{mouse_x - 210}+{mouse_y - 150}")

        # 文字输入
        tk.Label(edit_window, text="文字内容", bg="#ffffff",
                fg="#333333", font=("微软雅黑", 9, "bold")).pack(anchor="w", padx=15, pady=(15, 5))

        text_input = tk.Text(edit_window, height=4, bg="#f5f5f5",
                            font=("微软雅黑", 11), relief=tk.GROOVE, bd=1, wrap=tk.WORD)
        text_input.pack(fill=tk.X, padx=15, pady=5)
        text_input.insert("1.0", box.text)
        text_input.focus_set()
        text_input.tag_add("sel", "1.0", "end")

        # 快捷设置
        quick_frame = tk.Frame(edit_window, bg="#ffffff")
        quick_frame.pack(fill=tk.X, padx=15, pady=10)

        tk.Label(quick_frame, text="字号:", bg="#ffffff", font=("微软雅黑", 9)).pack(side=tk.LEFT)
        font_size_var = tk.StringVar(value=str(box.font_size))
        ttk.Combobox(quick_frame, textvariable=font_size_var, width=5,
                    values=["8", "10", "12", "14", "16", "18", "20", "24", "28", "32", "36", "48", "60", "72", "80", "100", "120", "150", "200"]).pack(side=tk.LEFT, padx=5)

        tk.Label(quick_frame, text="对齐:", bg="#ffffff", font=("微软雅黑", 9)).pack(side=tk.LEFT, padx=(15, 0))
        align_var = tk.StringVar(value=box.align)
        for val, txt in [("left", "左"), ("center", "中"), ("right", "右")]:
            tk.Radiobutton(quick_frame, text=txt, variable=align_var, value=val,
                          bg="#ffffff", font=("微软雅黑", 9)).pack(side=tk.LEFT)

        # 样式
        style_frame = tk.Frame(edit_window, bg="#ffffff")
        style_frame.pack(fill=tk.X, padx=15, pady=5)

        bold_var = tk.BooleanVar(value=box.bold)
        tk.Checkbutton(style_frame, text="加粗", variable=bold_var,
                      bg="#ffffff", font=("微软雅黑", 9)).pack(side=tk.LEFT)

        color_var = tk.StringVar(value=box.font_color)
        color_btn = tk.Button(style_frame, text="颜色", bg=box.font_color, width=6,
                             command=lambda: self._pick_color_for_editor(color_btn, color_var))
        color_btn.pack(side=tk.LEFT, padx=10)

        def auto_calc():
            text = text_input.get("1.0", tk.END).strip()
            if text:
                text_len = len(text)
                font_h = int(box.height * 0.7 * 72 / 96)
                font_w = int(box.width / text_len * 0.85 * 72 / 96) if text_len > 0 else font_h
                font_size_var.set(str(max(8, min(min(font_h, font_w), 72))))

        tk.Button(style_frame, text="自动字号", command=auto_calc,
                 bg="#9C27B0", fg="white", font=("微软雅黑", 9)).pack(side=tk.LEFT, padx=10)

        # 按钮
        btn_frame = tk.Frame(edit_window, bg="#ffffff")
        btn_frame.pack(fill=tk.X, padx=15, pady=15)

        def save():
            box.text = text_input.get("1.0", tk.END).strip()
            try:
                box.font_size = int(font_size_var.get())
            except:
                pass
            box.align = align_var.get()
            box.bold = bold_var.get()
            box.font_color = color_var.get()
            edit_window.destroy()
            self.refresh_canvas()
            self.update_listbox()
            self.update_property_panel()

        tk.Button(btn_frame, text="确定", command=save,
                 bg="#4CAF50", fg="white", font=("微软雅黑", 10),
                 width=10, cursor="hand2").pack(side=tk.LEFT, padx=5)
        tk.Button(btn_frame, text="取消", command=edit_window.destroy,
                 bg="#9E9E9E", fg="white", font=("微软雅黑", 10),
                 width=10, cursor="hand2").pack(side=tk.LEFT, padx=5)

        edit_window.bind("<Control-Return>", lambda e: save())
        edit_window.bind("<Escape>", lambda e: edit_window.destroy())

    def _pick_color_for_editor(self, btn, color_var):
        color = colorchooser.askcolor(color=color_var.get(), title="选择颜色")
        if color[1]:
            color_var.set(color[1])
            btn.config(bg=color[1])

    # ==================== 撤销/重做 ====================

    def save_state(self):
        state = [box.to_dict() for box in self.text_boxes]
        if self.history_index < len(self.history) - 1:
            self.history = self.history[:self.history_index + 1]
        self.history.append(state)
        if len(self.history) > self.max_history:
            self.history.pop(0)
        else:
            self.history_index += 1

    def undo(self):
        if self.history_index > 0:
            self.history_index -= 1
            self._restore_state(self.history[self.history_index])
            self.update_status("撤销 ✓")

    def redo(self):
        if self.history_index < len(self.history) - 1:
            self.history_index += 1
            self._restore_state(self.history[self.history_index])
            self.update_status("重做 ✓")

    def _restore_state(self, state):
        self.text_boxes = [TextBox.from_dict(data) for data in state]
        self.selected_box_index = -1
        self.selected_boxes = []
        self.refresh_canvas()
        self.update_listbox()
        self.mark_unsaved()

    # ==================== 框操作 ====================

    def delete_selected_box(self):
        if self.selected_box_index < 0:
            return
        self.save_state()
        del self.text_boxes[self.selected_box_index]
        self.selected_box_index = -1
        self.selected_boxes = []
        self.refresh_canvas()
        self.update_listbox()
        self.mark_unsaved()

    def clear_all_boxes(self):
        if messagebox.askyesno("确认", "确定清空所有文本框？"):
            self.save_state()
            self.text_boxes = []
            self.selected_box_index = -1
            self.selected_boxes = []
            self.refresh_canvas()
            self.update_listbox()

    def auto_font_size(self):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        if not box.text:
            return
        text_len = len(box.text)
        font_h = int(box.height * 0.7 * 72 / 96)
        font_w = int(box.width / text_len * 0.85 * 72 / 96) if text_len > 0 else font_h
        box.font_size = max(8, min(min(font_h, font_w), 72))
        self.fontsize_var.set(str(box.font_size))
        self.refresh_canvas()

    def auto_font_size_all(self):
        for box in self.text_boxes:
            if not box.text:
                continue
            text_len = len(box.text)
            font_h = int(box.height * 0.7 * 72 / 96)
            font_w = int(box.width / text_len * 0.85 * 72 / 96) if text_len > 0 else font_h
            box.font_size = max(8, min(min(font_h, font_w), 72))
        self.update_property_panel()
        self.refresh_canvas()
        self.update_status("已为当前页所有框计算字号 ✓")

    def align_boxes(self, align_type):
        if len(self.selected_boxes) < 2:
            self.update_status("请Ctrl+点击选中至少2个框")
            return

        self.save_state()
        boxes = [self.text_boxes[i] for i in self.selected_boxes]

        if align_type == "left":
            min_x = min(b.x for b in boxes)
            for b in boxes: b.x = min_x
        elif align_type == "right":
            max_right = max(b.x + b.width for b in boxes)
            for b in boxes: b.x = max_right - b.width
        elif align_type == "center_h":
            avg = sum(b.x + b.width / 2 for b in boxes) / len(boxes)
            for b in boxes: b.x = int(avg - b.width / 2)
        elif align_type == "top":
            min_y = min(b.y for b in boxes)
            for b in boxes: b.y = min_y
        elif align_type == "bottom":
            max_bottom = max(b.y + b.height for b in boxes)
            for b in boxes: b.y = max_bottom - b.height
        elif align_type == "center_v":
            avg = sum(b.y + b.height / 2 for b in boxes) / len(boxes)
            for b in boxes: b.y = int(avg - b.height / 2)

        self.refresh_canvas()
        self.update_status(f"已对齐 {len(self.selected_boxes)} 个框 ✓")

    def batch_offset(self, dx_dir, dy_dir):
        """批量位移选中的文本框

        Args:
            dx_dir: X方向（-1左, 0无, 1右）
            dy_dir: Y方向（-1上, 0无, 1下）
        """
        # 至少要有一个选中的框（包括主选中框）
        boxes_to_move = []
        if self.selected_boxes:
            boxes_to_move = self.selected_boxes
        elif self.selected_box_index >= 0:
            boxes_to_move = [self.selected_box_index]

        if not boxes_to_move:
            self.update_status("请先选中至少一个文本框")
            return

        # 获取像素值
        try:
            pixels = int(self.offset_px_var.get())
            if pixels <= 0:
                self.update_status("像素值必须大于0")
                return
        except ValueError:
            self.update_status("请输入有效的像素数值")
            return

        # 保存状态用于撤销
        self.save_state()

        # 计算实际偏移量
        dx = dx_dir * pixels
        dy = dy_dir * pixels

        # 移动所有选中的框
        for idx in boxes_to_move:
            if 0 <= idx < len(self.text_boxes):
                box = self.text_boxes[idx]
                box.x = max(0, box.x + dx)  # 不能移出边界
                box.y = max(0, box.y + dy)

        # 更新界面
        self.refresh_canvas()
        self.update_property_panel()
        self.mark_unsaved()

        # 提示信息
        direction = ""
        if dx_dir == -1:
            direction = "左"
        elif dx_dir == 1:
            direction = "右"
        elif dy_dir == -1:
            direction = "上"
        elif dy_dir == 1:
            direction = "下"

        self.update_status(f"已将 {len(boxes_to_move)} 个框向{direction}移动 {pixels} 像素 ✓")

    def apply_style_to_selected(self):
        if len(self.selected_boxes) < 1:
            self.update_status("请先Ctrl+点击选中框")
            return

        any_selected = (self.apply_fontsize_var.get() or self.apply_fontname_var.get() or
                       self.apply_color_var.get() or self.apply_bold_var.get() or
                       self.apply_italic_var.get() or self.apply_align_var.get())

        if not any_selected:
            self.update_status("请先勾选要应用的属性")
            return

        self.save_state()

        try:
            font_size = int(self.fontsize_var.get())
        except:
            font_size = 16

        for idx in self.selected_boxes:
            if 0 <= idx < len(self.text_boxes):
                box = self.text_boxes[idx]
                if self.apply_fontsize_var.get(): box.font_size = font_size
                if self.apply_fontname_var.get(): box.font_name = self.fontname_var.get()
                if self.apply_bold_var.get(): box.bold = self.bold_var.get()
                if self.apply_italic_var.get(): box.italic = self.italic_var.get()
                if self.apply_align_var.get(): box.align = self.align_var.get()
                if self.apply_color_var.get(): box.font_color = self.color_btn.cget("bg")

        self.refresh_canvas()
        self.update_status(f"已应用样式到 {len(self.selected_boxes)} 个框 ✓")

    # ==================== OCR ====================

    def _prepare_image_for_ocr(self, img_path, edit_scale=1.0):
        """准备OCR用的图片，如果图片过大则缩放"""
        MAX_SIDE = 3000  # 最大边长限制

        img = Image.open(img_path)
        w, h = img.size

        # 先应用编辑缩放
        if edit_scale < 1.0:
            w = int(w * edit_scale)
            h = int(h * edit_scale)
            img = img.resize((w, h), Image.Resampling.LANCZOS)

        # 如果还是太大，再缩放
        if max(w, h) <= MAX_SIDE:
            # 保存到临时文件
            temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
            temp_path = temp_file.name
            temp_file.close()
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            img.save(temp_path, quality=95)
            return temp_path, 1.0

        # 计算额外缩放比例
        extra_scale = MAX_SIDE / max(w, h)
        new_w = int(w * extra_scale)
        new_h = int(h * extra_scale)

        # 缩放图片
        resized_img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)

        # 保存到临时文件
        temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
        temp_path = temp_file.name
        temp_file.close()

        if resized_img.mode == 'RGBA':
            resized_img = resized_img.convert('RGB')
        resized_img.save(temp_path, quality=95)

        return temp_path, extra_scale

    def auto_detect_text_regions(self):
        if not self.original_image:
            self.update_status("请先加载图片")
            return
        if not self.ocr:
            self.update_status("OCR模型未加载")
            return

        if self.text_boxes:
            result = messagebox.askyesnocancel("提示", "是否清空现有框？\n是-清空  否-追加  取消-取消")
            if result is None:
                return
            elif result:
                self.text_boxes = []

        self.update_status("正在检测...")

        def detect():
            try:
                # 直接使用当前编辑图片，完全不缩放，保证坐标100%准确
                # PIL Image转为OpenCV格式
                img = np.array(self.original_image)
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

                img_h, img_w = img.shape[:2]

                # 保存临时文件用于OCR（不缩放！）
                temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                cv2.imwrite(temp_path, img)

                result = self.ocr.predict(temp_path)

                # 删除临时文件
                try:
                    os.remove(temp_path)
                except:
                    pass

                # 新版 PaddleOCR 返回 list，取第一个结果
                if not result or len(result) == 0:
                    self.root.after(0, lambda: self.update_status("未检测到文字"))
                    return

                ocr_result = result[0]
                dt_polys = ocr_result.get('dt_polys', [])
                rec_texts = ocr_result.get('rec_texts', [])

                if not dt_polys:
                    self.root.after(0, lambda: self.update_status("未检测到文字"))
                    return

                new_boxes = []
                for i, poly in enumerate(dt_polys):
                    x_coords = [p[0] for p in poly]
                    y_coords = [p[1] for p in poly]

                    # 完全使用OCR原始坐标，不做任何调整
                    x = int(min(x_coords))
                    y = int(min(y_coords))
                    w = int(max(x_coords) - min(x_coords))
                    h = int(max(y_coords) - min(y_coords))

                    if w < 10 or h < 10:
                        continue

                    box = TextBox(max(0, x), max(0, y), w, h)
                    if i < len(rec_texts):
                        box.text = rec_texts[i]
                    if box.text:
                        text_len = len(box.text)
                        font_h = int(h * 0.7 * 72 / 96)
                        font_w = int(w / text_len * 0.85 * 72 / 96) if text_len > 0 else font_h
                        box.font_size = max(8, min(min(font_h, font_w), 72))
                    new_boxes.append(box)

                new_boxes.sort(key=lambda b: (b.y // 30, b.x))
                self.text_boxes.extend(new_boxes)

                self.root.after(0, self.refresh_canvas)
                self.root.after(0, self.update_listbox)
                self.root.after(0, lambda: self.update_status(f"检测到 {len(new_boxes)} 个文字区域 ✓"))

            except Exception as e:
                self.root.after(0, lambda: self.update_status(f"检测失败: {e}"))

        threading.Thread(target=detect, daemon=True).start()

    def ocr_all_boxes(self):
        if not self.text_boxes or not self.ocr:
            return

        self.update_status("正在识别...")

        def ocr_task():
            # 使用当前编辑图片
            if self.original_image is None:
                self.root.after(0, lambda: self.update_status("无法读取图片"))
                return

            # PIL Image转为OpenCV格式
            img = np.array(self.original_image)
            img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

            img_h, img_w = img.shape[:2]

            for box in self.text_boxes:
                if box.text:
                    continue

                x, y, w, h = box.x, box.y, box.width, box.height
                expand_h, expand_w = int(h * 0.3), int(w * 0.1)

                crop_x = max(0, x - expand_w)
                crop_y = max(0, y - expand_h)
                crop_x2 = min(x + w + expand_w, img_w)
                crop_y2 = min(y + h + expand_h, img_h)

                cropped = img[crop_y:crop_y2, crop_x:crop_x2]

                temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                cv2.imwrite(temp_path, cropped)

                try:
                    result = self.ocr.predict(temp_path)
                    os.remove(temp_path)

                    if result and len(result) > 0:
                        ocr_result = result[0]
                        rec_texts = ocr_result.get('rec_texts', [])
                        if rec_texts:
                            box.text = ''.join(rec_texts)
                            if box.text:
                                text_len = len(box.text)
                                font_h = int(h * 0.7 * 72 / 96)
                                font_w = int(w / text_len * 0.85 * 72 / 96) if text_len > 0 else font_h
                                box.font_size = max(8, min(min(font_h, font_w), 72))
                except:
                    try:
                        os.remove(temp_path)
                    except:
                        pass

            self.root.after(0, self.refresh_canvas)
            self.root.after(0, self.update_listbox)
            self.root.after(0, self.update_property_panel)
            self.root.after(0, lambda: self.update_status("识别完成 ✓"))

        threading.Thread(target=ocr_task, daemon=True).start()

    def ocr_single_box(self):
        """OCR识别单个选中的文本框"""
        # 检查是否选中了文本框
        if self.selected_box_index < 0 or self.selected_box_index >= len(self.text_boxes):
            messagebox.showinfo("提示", "请先选中一个文本框")
            return

        # 检查OCR模型
        if not self.ocr:
            messagebox.showwarning("提示", "OCR模型正在加载中，请稍候...")
            return

        # 检查是否有原图
        if not self.original_img_path or not os.path.exists(self.original_img_path):
            messagebox.showerror("错误", "找不到原始图片")
            return

        box = self.text_boxes[self.selected_box_index]
        self.update_status(f"正在识别第 {self.selected_box_index + 1} 个文本框...")

        def ocr_task():
            try:
                # 使用编辑图片（self.original_image是PIL Image，需要转为cv2格式）
                # 这样坐标系统就对齐了，不会错位
                if self.original_image is None:
                    self.root.after(0, lambda: messagebox.showerror("错误", "无法读取图片"))
                    return

                # PIL Image转为OpenCV格式
                import numpy as np
                img = np.array(self.original_image)
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

                img_h, img_w = img.shape[:2]

                # 获取文本框区域（稍微扩大一点以提高识别率）
                x, y, w, h = box.x, box.y, box.width, box.height
                expand_h, expand_w = int(h * 0.3), int(w * 0.1)

                crop_x = max(0, x - expand_w)
                crop_y = max(0, y - expand_h)
                crop_x2 = min(x + w + expand_w, img_w)
                crop_y2 = min(y + h + expand_h, img_h)

                # 裁剪区域
                cropped = img[crop_y:crop_y2, crop_x:crop_x2]

                # 保存临时文件
                temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                cv2.imwrite(temp_path, cropped)

                # OCR识别
                try:
                    result = self.ocr.predict(temp_path)
                    os.remove(temp_path)

                    if result and len(result) > 0:
                        ocr_result = result[0]
                        rec_texts = ocr_result.get('rec_texts', [])

                        if rec_texts:
                            # 合并识别的文字
                            recognized_text = ''.join(rec_texts)

                            if recognized_text:
                                box.text = recognized_text

                                # 自动计算合适的字号
                                text_len = len(box.text)
                                font_h = int(h * 0.7 * 72 / 96)
                                font_w = int(w / text_len * 0.85 * 72 / 96) if text_len > 0 else font_h
                                box.font_size = max(8, min(min(font_h, font_w), 72))

                                # 更新界面
                                self.root.after(0, self.refresh_canvas)
                                self.root.after(0, self.update_listbox)
                                self.root.after(0, self.update_property_panel)
                                self.root.after(0, lambda: self.update_status(f"识别成功: {recognized_text[:20]}..."))
                                self.root.after(0, lambda: messagebox.showinfo("识别成功",
                                    f"识别结果：\n\n{recognized_text}\n\n"
                                    f"字号已自动调整为: {box.font_size}"))
                            else:
                                self.root.after(0, lambda: self.update_status("未识别到文字"))
                                self.root.after(0, lambda: messagebox.showwarning("识别结果", "未识别到文字"))
                        else:
                            self.root.after(0, lambda: self.update_status("未识别到文字"))
                            self.root.after(0, lambda: messagebox.showwarning("识别结果", "未识别到文字"))
                    else:
                        self.root.after(0, lambda: self.update_status("识别失败"))
                        self.root.after(0, lambda: messagebox.showwarning("识别结果", "未识别到文字"))

                except Exception as e:
                    try:
                        os.remove(temp_path)
                    except:
                        pass
                    self.root.after(0, lambda: messagebox.showerror("错误", f"OCR识别出错:\n{str(e)}"))
                    self.root.after(0, lambda: self.update_status("识别失败"))

            except Exception as e:
                import traceback
                traceback.print_exc()
                self.root.after(0, lambda: messagebox.showerror("错误", f"识别出错:\n{str(e)}"))
                self.root.after(0, lambda: self.update_status("识别失败"))

        threading.Thread(target=ocr_task, daemon=True).start()

    # ==================== 批量操作 ====================

    def auto_detect_all_pages(self):
        if not self.pages or not self.ocr:
            self.update_status("请先导入图片")
            return

        self.save_current_page()

        def detect_all():
            total = len(self.pages)
            for i, page in enumerate(self.pages):
                self.root.after(0, lambda idx=i: self.update_status(f"检测第 {idx+1}/{total} 页..."))

                try:
                    # 直接使用该页的编辑图片，完全不缩放
                    page_img = page["image"]
                    img = np.array(page_img)
                    img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

                    # 保存临时文件（不缩放！）
                    temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                    temp_path = temp_file.name
                    temp_file.close()
                    cv2.imwrite(temp_path, img)

                    result = self.ocr.predict(temp_path)

                    # 删除临时文件
                    try:
                        os.remove(temp_path)
                    except:
                        pass

                    if not result or len(result) == 0:
                        continue

                    ocr_result = result[0]
                    dt_polys = ocr_result.get('dt_polys', [])
                    rec_texts = ocr_result.get('rec_texts', [])

                    if not dt_polys:
                        continue

                    new_boxes = []
                    for j, poly in enumerate(dt_polys):
                        x_coords = [p[0] for p in poly]
                        y_coords = [p[1] for p in poly]

                        # 完全使用OCR原始坐标，不做任何调整
                        x = int(min(x_coords))
                        y = int(min(y_coords))
                        w = int(max(x_coords) - min(x_coords))
                        h = int(max(y_coords) - min(y_coords))

                        if w < 10 or h < 10:
                            continue

                        box_data = {
                            "x": max(0, x), "y": max(0, y), "width": w, "height": h,
                            "text": rec_texts[j] if j < len(rec_texts) else "",
                            "font_size": 16, "font_name": "微软雅黑", "font_color": "#000000",
                            "bold": False, "italic": False, "align": "left"
                        }

                        if box_data["text"]:
                            text_len = len(box_data["text"])
                            font_h = int(h * 0.7 * 72 / 96)
                            font_w = int(w / text_len * 0.85 * 72 / 96) if text_len > 0 else font_h
                            box_data["font_size"] = max(8, min(min(font_h, font_w), 72))

                        new_boxes.append(box_data)

                    new_boxes.sort(key=lambda b: (b["y"] // 30, b["x"]))
                    page["text_boxes"] = new_boxes

                except Exception as e:
                    print(f"第 {i+1} 页检测失败: {e}")

            self.root.after(0, self.load_current_page)
            self.root.after(0, lambda: self.update_status(f"全部检测完成！共 {total} 页 ✓"))

        threading.Thread(target=detect_all, daemon=True).start()

    def ocr_all_pages(self):
        if not self.pages or not self.ocr:
            return

        self.save_current_page()

        def ocr_all():
            total = len(self.pages)
            for i, page in enumerate(self.pages):
                self.root.after(0, lambda idx=i: self.update_status(f"识别第 {idx+1}/{total} 页..."))

                boxes = page.get("text_boxes", [])
                if not boxes:
                    continue

                # 使用该页的编辑图片
                page_img = page["image"]
                img = np.array(page_img)
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

                img_h, img_w = img.shape[:2]

                for box_data in boxes:
                    if box_data.get("text"):
                        continue

                    x, y, w, h = box_data["x"], box_data["y"], box_data["width"], box_data["height"]
                    expand_h, expand_w = int(h * 0.3), int(w * 0.1)

                    crop_x = max(0, x - expand_w)
                    crop_y = max(0, y - expand_h)
                    crop_x2 = min(x + w + expand_w, img_w)
                    crop_y2 = min(y + h + expand_h, img_h)

                    cropped = img[crop_y:crop_y2, crop_x:crop_x2]

                    temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                    temp_path = temp_file.name
                    temp_file.close()
                    cv2.imwrite(temp_path, cropped)

                    try:
                        result = self.ocr.predict(temp_path)
                        os.remove(temp_path)

                        if result and len(result) > 0:
                            ocr_result = result[0]
                            rec_texts = ocr_result.get('rec_texts', [])
                            if rec_texts:
                                box_data["text"] = ''.join(rec_texts)
                                if box_data["text"]:
                                    text_len = len(box_data["text"])
                                    font_h = int(h * 0.7 * 72 / 96)
                                    font_w = int(w / text_len * 0.85 * 72 / 96) if text_len > 0 else font_h
                                    box_data["font_size"] = max(8, min(min(font_h, font_w), 72))
                    except:
                        try:
                            os.remove(temp_path)
                        except:
                            pass

            self.root.after(0, self.load_current_page)
            self.root.after(0, lambda: self.update_status(f"全部识别完成！共 {total} 页 ✓"))

        threading.Thread(target=ocr_all, daemon=True).start()

    def auto_font_size_all_pages(self):
        if not self.pages:
            return

        self.save_current_page()

        for page in self.pages:
            for box_data in page.get("text_boxes", []):
                if not box_data.get("text"):
                    continue
                text_len = len(box_data["text"])
                h, w = box_data["height"], box_data["width"]
                font_h = int(h * 0.7 * 72 / 96)
                font_w = int(w / text_len * 0.85 * 72 / 96) if text_len > 0 else font_h
                box_data["font_size"] = max(8, min(min(font_h, font_w), 72))

        self.load_current_page()
        self.update_status(f"全部 {len(self.pages)} 页字号已调整 ✓")

    # ==================== 项目保存/加载 ====================

    def save_project(self):
        self.save_current_page()

        file_path = filedialog.asksaveasfilename(
            defaultextension=".json",
            filetypes=[("JSON文件", "*.json")]
        )
        if not file_path:
            return

        pages_data = []
        for page in self.pages:
            pages_data.append({
                "original_path": page["original_path"],
                "original_size": page.get("original_size", page["image"].size),
                "edit_scale": page.get("edit_scale", 1.0),
                "bg_path": page.get("bg_path"),
                "bg_original_path": page.get("bg_original_path"),
                "text_boxes": page.get("text_boxes", [])
            })

        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump({"version": 3, "pages": pages_data, "current_page": self.current_page_index},
                     f, ensure_ascii=False, indent=2)

        self.update_status(f"项目已保存: {len(self.pages)} 页 ✓")
        self.mark_saved()

    def load_project(self):
        file_path = filedialog.askopenfilename(filetypes=[("JSON文件", "*.json")])
        if not file_path:
            return

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)

            self.pages = []
            for page_info in data.get("pages", []):
                if os.path.exists(page_info["original_path"]):
                    original_img = Image.open(page_info["original_path"])
                    original_size = page_info.get("original_size", original_img.size)

                    # 缩放图片用于编辑
                    edit_img, edit_scale = self._resize_image_for_edit(original_img)

                    self.pages.append({
                        "original_path": page_info["original_path"],
                        "original_size": original_size,
                        "edit_scale": edit_scale,
                        "bg_path": page_info.get("bg_path"),
                        "bg_original_path": page_info.get("bg_original_path"),
                        "image": edit_img,
                        "text_boxes": page_info.get("text_boxes", [])
                    })

            self.current_page_index = min(data.get("current_page", 0), len(self.pages) - 1 if self.pages else 0)

            if self.pages:
                self.load_current_page()
                self.update_page_label()
                self.update_thumbnails()
                self.placeholder_label.place_forget()

            self.update_status(f"已加载 {len(self.pages)} 页项目 ✓")
            self.mark_saved()

        except Exception as e:
            messagebox.showerror("错误", f"加载失败: {e}")

    # ==================== PPT生成 ====================

    def generate_multi_page_ppt(self):
        if not self.pages:
            self.update_status("请先导入图片")
            return

        self.save_current_page()

        save_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint文件", "*.pptx")]
        )
        if not save_path:
            return

        try:
            prs = Presentation()

            for page_idx, page in enumerate(self.pages):
                self.update_status(f"生成第 {page_idx+1}/{len(self.pages)} 页...")

                # 直接使用编辑图片的尺寸（已缩放）
                img_w, img_h = page["image"].size

                if page_idx == 0:
                    prs.slide_width = Px(img_w)
                    prs.slide_height = Px(img_h)

                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # 使用缩放后的背景图（bg_path已经是缩放后的）
                bg_path = page.get("bg_path") or page["original_path"]
                # 如果使用的是原图路径但原图很大，需要使用编辑图片
                if bg_path == page["original_path"] and page.get("edit_scale", 1.0) < 1.0:
                    # 保存编辑图片到临时文件作为背景
                    temp_bg = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                    temp_bg_path = temp_bg.name
                    temp_bg.close()
                    edit_img = page["image"]
                    if edit_img.mode == 'RGBA':
                        edit_img = edit_img.convert('RGB')
                    edit_img.save(temp_bg_path, quality=95)
                    bg_path = temp_bg_path

                slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

                for box_data in page.get("text_boxes", []):
                    if not box_data.get("text"):
                        continue

                    # 直接使用编辑坐标（不需要转换）
                    x = box_data["x"]
                    y = box_data["y"]
                    w = box_data["width"]
                    h = box_data["height"]
                    font_size = box_data.get("font_size", 16)

                    textbox = slide.shapes.add_textbox(
                        Px(x), Px(y), Px(w), Px(h)
                    )
                    tf = textbox.text_frame
                    tf.word_wrap = False
                    tf.margin_left = Px(2)
                    tf.margin_right = Px(2)
                    tf.margin_top = Px(1)
                    tf.margin_bottom = Px(1)
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                    p = tf.paragraphs[0]
                    p.text = box_data["text"]

                    align = box_data.get("align", "left")
                    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)

                    if p.runs:
                        run = p.runs[0]
                        run.font.size = Pt(font_size)
                        run.font.name = box_data.get("font_name", "微软雅黑")
                        run.font.bold = box_data.get("bold", False)
                        run.font.italic = box_data.get("italic", False)

                        color_hex = box_data.get("font_color", "#000000").lstrip('#')
                        run.font.color.rgb = RGBColor(int(color_hex[0:2], 16),
                                                      int(color_hex[2:4], 16),
                                                      int(color_hex[4:6], 16))

                    p.line_spacing = 1.0
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)

            prs.save(save_path)
            messagebox.showinfo("成功", f"PPT已保存！\n共 {len(self.pages)} 页\n{save_path}")
            self.update_status(f"PPT生成成功！共 {len(self.pages)} 页 ✓")

        except Exception as e:
            messagebox.showerror("失败", f"生成失败: {e}")

    # ==================== 设置对话框 ====================

    def show_settings_dialog(self):
        """显示设置对话框"""
        dialog = tk.Toplevel(self.root)
        dialog.title("设置")
        dialog.geometry("680x700")  # 增大对话框尺寸
        dialog.configure(bg=COLOR_WHITE)
        dialog.transient(self.root)
        dialog.grab_set()

        # 居中显示
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() - 680) // 2
        y = (dialog.winfo_screenheight() - 700) // 2
        dialog.geometry(f"+{x}+{y}")

        # 标题
        title_frame = tk.Frame(dialog, bg=COLOR_THEME, height=40)
        title_frame.pack(fill=tk.X, side=tk.TOP)
        title_frame.pack_propagate(False)
        tk.Label(title_frame, text="  OCR模型设置", bg=COLOR_THEME, fg="white",
                font=(FONT_FAMILY, 12, "bold")).pack(side=tk.LEFT, pady=8)

        # 按钮区 - 固定在底部
        btn_frame = tk.Frame(dialog, bg=COLOR_WHITE, pady=15)
        btn_frame.pack(fill=tk.X, side=tk.BOTTOM)

        tk.Button(btn_frame, text="保存并加载OCR", command=lambda: self._save_settings(dialog),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 11, "bold"),
                 padx=30, pady=8, cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=20)

        tk.Button(btn_frame, text="取消", command=dialog.destroy,
                 bg="#9E9E9E", fg="white", font=(FONT_FAMILY, 11),
                 padx=30, pady=8, cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT)

        # 分隔线
        tk.Frame(dialog, bg="#ddd", height=1).pack(fill=tk.X, side=tk.BOTTOM)

        # 可滚动内容区 - 放在中间
        content_container = tk.Frame(dialog, bg=COLOR_WHITE)
        content_container.pack(fill=tk.BOTH, expand=True, side=tk.TOP)

        canvas = tk.Canvas(content_container, bg=COLOR_WHITE, highlightthickness=0)
        scrollbar = tk.Scrollbar(content_container, orient=tk.VERTICAL, command=canvas.yview)

        content = tk.Frame(canvas, bg=COLOR_WHITE, padx=20, pady=15)

        canvas.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        canvas_window = canvas.create_window((0, 0), window=content, anchor=tk.NW)

        # 更新滚动区域
        def on_frame_configure(event):
            canvas.configure(scrollregion=canvas.bbox("all"))

        content.bind("<Configure>", on_frame_configure)

        # 调整canvas窗口宽度
        def on_canvas_configure(event):
            canvas.itemconfig(canvas_window, width=event.width)

        canvas.bind("<Configure>", on_canvas_configure)

        # 鼠标滚轮支持
        def on_mousewheel(event):
            canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")

        canvas.bind_all("<MouseWheel>", on_mousewheel)

        # === 方式1: 指定已有模型目录 ===
        tk.Label(content, text="方式1: 指定已有模型目录", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w")
        tk.Label(content, text="如果已有模型文件，直接选择模型所在目录",
                bg=COLOR_WHITE, fg="#666", font=(FONT_FAMILY, 9)).pack(anchor="w", pady=(0, 5))

        path_frame = tk.Frame(content, bg=COLOR_WHITE)
        path_frame.pack(fill=tk.X, pady=5)

        self.model_dir_var = tk.StringVar(value=self.config.get("model_dir", ""))
        path_entry = tk.Entry(path_frame, textvariable=self.model_dir_var,
                             font=(FONT_FAMILY, 10), width=45)
        path_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)

        browse_btn = tk.Button(path_frame, text="浏览...", command=self._browse_model_dir,
                              bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9),
                              padx=10, cursor="hand2", relief=tk.FLAT)
        browse_btn.pack(side=tk.LEFT, padx=(10, 0))

        # === 方式2: 下载模型到指定目录 ===
        tk.Frame(content, bg="#ddd", height=1).pack(fill=tk.X, pady=15)

        tk.Label(content, text="方式2: 下载模型到指定目录", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w")
        tk.Label(content, text="如果没有模型，选择一个目录后点击下载（需要联网，约200MB）",
                bg=COLOR_WHITE, fg="#666", font=(FONT_FAMILY, 9)).pack(anchor="w", pady=(0, 5))

        download_frame = tk.Frame(content, bg=COLOR_WHITE)
        download_frame.pack(fill=tk.X, pady=5)

        self.download_dir_var = tk.StringVar(value=os.path.join(get_base_dir(), ".paddlex", "official_models"))
        download_entry = tk.Entry(download_frame, textvariable=self.download_dir_var,
                                 font=(FONT_FAMILY, 10), width=45)
        download_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)

        browse_download_btn = tk.Button(download_frame, text="浏览...",
                                       command=lambda: self._browse_download_dir(),
                                       bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9),
                                       padx=10, cursor="hand2", relief=tk.FLAT)
        browse_download_btn.pack(side=tk.LEFT, padx=(10, 0))

        # 下载按钮和进度
        download_btn_frame = tk.Frame(content, bg=COLOR_WHITE)
        download_btn_frame.pack(fill=tk.X, pady=10)

        self.download_btn = tk.Button(download_btn_frame, text="下载模型",
                                     command=lambda: self._download_models(dialog),
                                     bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"),
                                     padx=20, pady=5, cursor="hand2", relief=tk.FLAT)
        self.download_btn.pack(side=tk.LEFT)

        self.download_status_label = tk.Label(download_btn_frame, text="", bg=COLOR_WHITE,
                                             fg="#666", font=(FONT_FAMILY, 9))
        self.download_status_label.pack(side=tk.LEFT, padx=15)

        # 进度条
        progress_frame = tk.Frame(content, bg=COLOR_WHITE)
        progress_frame.pack(fill=tk.X, pady=5)

        self.download_progress = ttk.Progressbar(progress_frame, length=400, mode='determinate')
        self.download_progress.pack(fill=tk.X)

        self.download_detail_label = tk.Label(progress_frame, text="", bg=COLOR_WHITE,
                                              fg="#999", font=(FONT_FAMILY, 8))
        self.download_detail_label.pack(anchor="w")

        # === 设备选择 ===
        tk.Frame(content, bg="#ddd", height=1).pack(fill=tk.X, pady=15)

        tk.Label(content, text="设备选择", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w")
        tk.Label(content, text="选择OCR运行的设备（GPU需要安装PaddlePaddle-GPU版本）",
                bg=COLOR_WHITE, fg="#666", font=(FONT_FAMILY, 9)).pack(anchor="w", pady=(0, 5))

        device_frame = tk.Frame(content, bg=COLOR_WHITE)
        device_frame.pack(fill=tk.X, pady=5)

        self.device_var = tk.StringVar(value=self.config.get("ocr_device", "cpu"))

        tk.Radiobutton(device_frame, text="CPU - 兼容性好，适合所有电脑",
                      variable=self.device_var, value="cpu",
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 10)).pack(anchor="w", pady=3)
        tk.Radiobutton(device_frame, text="GPU - 速度快，需要NVIDIA显卡",
                      variable=self.device_var, value="gpu",
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 10)).pack(anchor="w", pady=3)

        # 提示信息
        tk.Label(device_frame,
                text="提示：使用GPU需要先安装 paddlepaddle-gpu\n如未安装，请运行：pip uninstall paddlepaddle && pip install paddlepaddle-gpu",
                bg=COLOR_WHITE, fg="#999", font=(FONT_FAMILY, 8), justify=tk.LEFT).pack(anchor="w", pady=(5, 0))

        # 模型状态显示
        tk.Frame(content, bg="#ddd", height=1).pack(fill=tk.X, pady=10)

        status_frame = tk.Frame(content, bg="#f5f5f5", padx=10, pady=10)
        status_frame.pack(fill=tk.X)

        self.model_status_label = tk.Label(status_frame, text="", bg="#f5f5f5",
                                           font=(FONT_FAMILY, 9), justify=tk.LEFT)
        self.model_status_label.pack(anchor="w")

        self._check_model_status()

        # 绑定路径变化事件
        self.model_dir_var.trace_add("write", lambda *args: self._check_model_status())

    def _browse_model_dir(self):
        """浏览选择模型目录"""
        current_dir = self.model_dir_var.get()
        if not os.path.exists(current_dir):
            current_dir = get_base_dir()

        dir_path = filedialog.askdirectory(
            title="选择OCR模型目录（包含 PP-OCRv5_server_det 等文件夹）",
            initialdir=current_dir
        )
        if dir_path:
            self.model_dir_var.set(dir_path)

    def _browse_download_dir(self):
        """浏览选择下载目录"""
        current_dir = self.download_dir_var.get()
        if not os.path.exists(current_dir):
            current_dir = get_base_dir()

        dir_path = filedialog.askdirectory(
            title="选择模型下载目录",
            initialdir=current_dir
        )
        if dir_path:
            self.download_dir_var.set(dir_path)

    def _download_models(self, dialog):
        """下载OCR模型 - 使用直接URL下载"""
        download_dir = self.download_dir_var.get()

        if not download_dir:
            messagebox.showwarning("警告", "请先选择下载目录！")
            return

        # 创建目录
        os.makedirs(download_dir, exist_ok=True)

        # 禁用下载按钮
        self.download_btn.config(state=tk.DISABLED, text="下载中...")
        self.download_status_label.config(text="正在准备下载...")
        self.download_progress['value'] = 0

        # 需要下载的模型列表
        models_to_download = [
            ("PP-OCRv5_server_det", "文字检测模型", "PP-OCRv5_server_det_infer.tar"),
            ("PP-OCRv5_server_rec", "文字识别模型", "PP-OCRv5_server_rec_infer.tar"),
            ("PP-LCNet_x1_0_doc_ori", "文档方向分类", "PP-LCNet_x1_0_doc_ori_infer.tar"),
            ("PP-LCNet_x1_0_textline_ori", "文本行方向", "PP-LCNet_x1_0_textline_ori_infer.tar"),
            ("UVDoc", "文档矫正", "UVDoc_infer.tar"),
        ]

        base_url = "https://paddle-model-ecology.bj.bcebos.com/paddlex/official_inference_model/paddle3.0.0"

        def download_task():
            import urllib.request
            import tarfile

            total_models = len(models_to_download)
            downloaded = 0

            for model_name, desc, tar_file in models_to_download:
                model_path = os.path.join(download_dir, model_name)

                # 如果模型已存在，跳过
                if os.path.exists(model_path):
                    downloaded += 1
                    progress = int((downloaded / total_models) * 100)
                    dialog.after(0, lambda p=progress, d=desc: self._update_download_progress(p, f"{d} 已存在，跳过"))
                    continue

                url = f"{base_url}/{tar_file}"
                tar_path = os.path.join(download_dir, tar_file)

                try:
                    # 更新状态
                    dialog.after(0, lambda d=desc: self.download_status_label.config(text=f"正在下载: {d}"))
                    dialog.after(0, lambda d=desc: self.download_detail_label.config(text=f"从 {url}"))

                    # 下载文件（带进度）
                    def reporthook(block_num, block_size, total_size):
                        if total_size > 0:
                            downloaded_size = block_num * block_size
                            percent = min(int((downloaded_size / total_size) * 100), 100)
                            size_mb = downloaded_size / (1024 * 1024)
                            total_mb = total_size / (1024 * 1024)
                            # 计算总进度
                            model_progress = downloaded / total_models
                            file_progress = (downloaded_size / total_size) / total_models
                            overall = int((model_progress + file_progress) * 100)
                            dialog.after(0, lambda o=overall, s=size_mb, t=total_mb:
                                self._update_download_progress(o, f"下载中: {s:.1f}MB / {t:.1f}MB"))

                    urllib.request.urlretrieve(url, tar_path, reporthook)

                    # 解压
                    dialog.after(0, lambda d=desc: self.download_status_label.config(text=f"正在解压: {d}"))

                    with tarfile.open(tar_path, 'r:*') as tar:
                        tar.extractall(download_dir)

                    # 删除tar文件
                    os.remove(tar_path)

                    # 重命名文件夹（去掉_infer后缀）
                    infer_path = os.path.join(download_dir, f"{model_name}_infer")
                    if os.path.exists(infer_path) and not os.path.exists(model_path):
                        os.rename(infer_path, model_path)

                    downloaded += 1
                    progress = int((downloaded / total_models) * 100)
                    dialog.after(0, lambda p=progress, d=desc: self._update_download_progress(p, f"{d} 下载完成"))

                except Exception as e:
                    dialog.after(0, lambda d=desc, err=str(e):
                        self.download_status_label.config(text=f"{d} 下载失败: {err[:50]}"))
                    # 清理可能的残留文件
                    if os.path.exists(tar_path):
                        try:
                            os.remove(tar_path)
                        except:
                            pass

            # 下载完成
            dialog.after(0, lambda: self._download_complete(download_dir, dialog))

        threading.Thread(target=download_task, daemon=True).start()

    def _update_download_progress(self, progress, detail):
        """更新下载进度"""
        self.download_progress['value'] = progress
        self.download_detail_label.config(text=detail)

    def _download_complete(self, download_dir, dialog):
        """下载完成处理"""
        self.download_btn.config(state=tk.NORMAL, text="下载模型")
        self.download_progress['value'] = 100
        self.download_status_label.config(text="下载完成！")
        self.download_detail_label.config(text="")

        # 设置模型目录
        self.model_dir_var.set(download_dir)
        self._check_model_status()

        messagebox.showinfo("成功",
            f"模型下载完成！\n\n下载目录:\n{download_dir}\n\n已自动设置为模型目录，点击'保存并加载OCR'即可使用。")

    def _check_model_status(self):
        """检查模型状态"""
        model_dir = self.model_dir_var.get()

        required_models = [
            ("PP-OCRv5_server_det", "文字检测模型"),
            ("PP-OCRv5_server_rec", "文字识别模型"),
        ]
        optional_models = [
            ("PP-LCNet_x1_0_doc_ori", "文档方向分类"),
            ("PP-LCNet_x1_0_textline_ori", "文本行方向"),
            ("UVDoc", "文档矫正"),
        ]

        status_lines = []

        if not model_dir:
            status_lines.append("请选择或下载模型目录")
        elif not os.path.exists(model_dir):
            status_lines.append("目录不存在，请选择有效目录或下载模型")
        else:
            all_required = True
            for model_name, desc in required_models:
                model_path = os.path.join(model_dir, model_name)
                if os.path.exists(model_path):
                    status_lines.append(f"[OK] {desc} ({model_name})")
                else:
                    status_lines.append(f"[X] {desc} ({model_name}) - 缺失!")
                    all_required = False

            for model_name, desc in optional_models:
                model_path = os.path.join(model_dir, model_name)
                if os.path.exists(model_path):
                    status_lines.append(f"[OK] {desc} ({model_name})")
                else:
                    status_lines.append(f"[  ] {desc} ({model_name}) - 可选")

            if all_required:
                status_lines.insert(0, "当前模型状态: 可用\n")
            else:
                status_lines.insert(0, "当前模型状态: 缺少必需模型!\n")

        self.model_status_label.config(text="\n".join(status_lines))

    def _save_settings(self, dialog):
        """保存设置并重新加载OCR"""
        new_model_dir = self.model_dir_var.get()
        new_device = self.device_var.get()  # 获取设备选择

        if not new_model_dir:
            messagebox.showwarning("警告", "请先选择模型目录！")
            return

        # 检查必需模型是否存在
        det_model = os.path.join(new_model_dir, "PP-OCRv5_server_det")
        rec_model = os.path.join(new_model_dir, "PP-OCRv5_server_rec")

        if not os.path.exists(det_model) or not os.path.exists(rec_model):
            result = messagebox.askyesno("警告",
                "模型目录缺少必需的模型文件！\n\n"
                "需要:\n- PP-OCRv5_server_det\n- PP-OCRv5_server_rec\n\n"
                "是否仍然保存？（OCR功能将无法使用）")
            if not result:
                return

        # 保存配置
        self.config["model_dir"] = new_model_dir
        self.config["ocr_device"] = new_device  # 保存设备选择
        save_config(self.config)

        # 关闭对话框
        dialog.destroy()

        # 重新加载OCR
        self.ocr = None
        device_name = "GPU" if new_device == "gpu" else "CPU"
        self.update_status(f"正在使用 {device_name} 加载OCR模型...")
        threading.Thread(target=self.init_ocr, daemon=True).start()

        messagebox.showinfo("成功",
            f"设置已保存！\n\n"
            f"模型目录:\n{new_model_dir}\n\n"
            f"运行设备: {device_name}\n\n"
            f"OCR模型正在后台加载...")



    # ==================== 新增功能：全选和复制粘贴 ====================

    def select_all_boxes(self):
        """全选当前页所有文本框"""
        if not self.text_boxes:
            self.update_status("当前页没有文本框")
            return

        # 选中所有框
        self.selected_boxes = list(range(len(self.text_boxes)))
        self.selected_box_index = 0 if self.text_boxes else -1

        # 刷新界面
        self.refresh_canvas()
        self.update_property_panel()

        # 更新列表框选择
        self.box_listbox.selection_clear(0, tk.END)
        for idx in self.selected_boxes:
            self.box_listbox.selection_set(idx)

        self.update_status(f"已选中当前页所有 {len(self.text_boxes)} 个文本框 ✓")

    def copy_boxes(self):
        """复制选中的文本框"""
        if not self.selected_boxes:
            self.update_status("请先选中要复制的文本框")
            return

        self.clipboard_boxes = []
        for idx in self.selected_boxes:
            if 0 <= idx < len(self.text_boxes):
                self.clipboard_boxes.append(self.text_boxes[idx].copy())

        self.update_status(f"已复制 {len(self.clipboard_boxes)} 个文本框")

    def paste_boxes(self):
        """粘贴文本框"""
        if not self.clipboard_boxes:
            self.update_status("剪贴板为空")
            return

        self.save_state()

        offset = 20
        new_boxes = []
        for box in self.clipboard_boxes:
            new_box = box.copy()
            new_box.x += offset
            new_box.y += offset
            self.text_boxes.append(new_box)
            new_boxes.append(new_box)

        start_idx = len(self.text_boxes) - len(new_boxes)
        self.selected_boxes = list(range(start_idx, len(self.text_boxes)))
        self.selected_box_index = self.selected_boxes[0] if self.selected_boxes else -1

        self.refresh_canvas()
        self.update_listbox()
        self.mark_unsaved()
        self.mark_unsaved()
        self.update_status(f"已粘贴 {len(new_boxes)} 个文本框")

    def move_box_by_key(self, dx, dy):
        """使用方向键移动文本框"""
        if self.selected_box_index < 0:
            return

        box = self.text_boxes[self.selected_box_index]
        box.x = max(0, box.x + dx)
        box.y = max(0, box.y + dy)

        self.refresh_canvas()
        self.update_property_panel()
        self.mark_unsaved()

    # ==================== 新增功能：完整对齐工具 ====================

    def show_align_dialog(self):
        """显示对齐工具对话框"""
        if len(self.selected_boxes) < 2:
            messagebox.showinfo("提示", "请先使用Ctrl+点击选中至少2个文本框")
            return

        dialog = tk.Toplevel(self.root)
        dialog.title("对齐与分布工具")
        dialog.geometry("450x550")
        dialog.configure(bg=COLOR_WHITE)
        dialog.transient(self.root)
        dialog.grab_set()

        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() - 450) // 2
        y = (dialog.winfo_screenheight() - 550) // 2
        dialog.geometry(f"+{x}+{y}")

        title_frame = tk.Frame(dialog, bg=COLOR_THEME, height=40)
        title_frame.pack(fill=tk.X)
        title_frame.pack_propagate(False)
        tk.Label(title_frame, text=f"  对齐与分布 - 已选中 {len(self.selected_boxes)} 个框",
                bg=COLOR_THEME, fg="white",
                font=(FONT_FAMILY, 11, "bold")).pack(side=tk.LEFT, pady=8)

        content = tk.Frame(dialog, bg=COLOR_WHITE, padx=20, pady=15)
        content.pack(fill=tk.BOTH, expand=True)

        # 水平对齐
        tk.Label(content, text="水平对齐", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(5, 5))

        h_frame = tk.Frame(content, bg=COLOR_WHITE)
        h_frame.pack(fill=tk.X, pady=5)

        tk.Button(h_frame, text="左对齐", command=lambda: self.align_boxes("left"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(h_frame, text="水平居中", command=lambda: self.align_boxes("center_h"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(h_frame, text="右对齐", command=lambda: self.align_boxes("right"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 垂直对齐
        tk.Label(content, text="垂直对齐", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(15, 5))

        v_frame = tk.Frame(content, bg=COLOR_WHITE)
        v_frame.pack(fill=tk.X, pady=5)

        tk.Button(v_frame, text="顶对齐", command=lambda: self.align_boxes("top"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(v_frame, text="垂直居中", command=lambda: self.align_boxes("center_v"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(v_frame, text="底对齐", command=lambda: self.align_boxes("bottom"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 分布
        tk.Label(content, text="均匀分布 (需要3个或以上)", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(15, 5))

        dist_frame = tk.Frame(content, bg=COLOR_WHITE)
        dist_frame.pack(fill=tk.X, pady=5)

        tk.Button(dist_frame, text="水平等间距", command=lambda: self.distribute_boxes("horizontal"),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 9), width=15,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(dist_frame, text="垂直等间距", command=lambda: self.distribute_boxes("vertical"),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 9), width=15,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 尺寸统一
        tk.Label(content, text="尺寸统一 (以第一个选中框为基准)", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(15, 5))

        size_frame = tk.Frame(content, bg=COLOR_WHITE)
        size_frame.pack(fill=tk.X, pady=5)

        tk.Button(size_frame, text="统一宽度", command=lambda: self.unify_size("width"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(size_frame, text="统一高度", command=lambda: self.unify_size("height"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(size_frame, text="统一大小", command=lambda: self.unify_size("both"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # 对齐到画布
        tk.Label(content, text="对齐到画布", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(15, 5))

        canvas_frame = tk.Frame(content, bg=COLOR_WHITE)
        canvas_frame.pack(fill=tk.X, pady=5)

        tk.Button(canvas_frame, text="画布水平居中", command=lambda: self.align_to_canvas("h"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 9), width=15,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(canvas_frame, text="画布垂直居中", command=lambda: self.align_to_canvas("v"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 9), width=15,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        canvas_frame2 = tk.Frame(content, bg=COLOR_WHITE)
        canvas_frame2.pack(fill=tk.X, pady=5)

        tk.Button(canvas_frame2, text="画布完全居中", command=lambda: self.align_to_canvas("center"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 9), width=32,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        tk.Frame(content, bg="#ddd", height=1).pack(fill=tk.X, pady=15)
        tk.Button(content, text="关闭", command=dialog.destroy,
                 bg=COLOR_GRAY, fg="white", font=(FONT_FAMILY, 10),
                 width=15, cursor="hand2", relief=tk.FLAT).pack()

    def distribute_boxes(self, direction):
        """均匀分布文本框"""
        if len(self.selected_boxes) < 3:
            messagebox.showinfo("提示", "均匀分布需要至少选中3个文本框")
            return

        self.save_state()
        boxes = [self.text_boxes[i] for i in self.selected_boxes]

        if direction == "horizontal":
            boxes.sort(key=lambda b: b.x)
            first = boxes[0]
            last = boxes[-1]

            total_width = sum(b.width for b in boxes)
            total_space = (last.x + last.width) - first.x - total_width
            gap = total_space / (len(boxes) - 1) if len(boxes) > 1 else 0

            current_x = first.x + first.width
            for box in boxes[1:-1]:
                box.x = int(current_x + gap)
                current_x = box.x + box.width

        elif direction == "vertical":
            boxes.sort(key=lambda b: b.y)
            first = boxes[0]
            last = boxes[-1]

            total_height = sum(b.height for b in boxes)
            total_space = (last.y + last.height) - first.y - total_height
            gap = total_space / (len(boxes) - 1) if len(boxes) > 1 else 0

            current_y = first.y + first.height
            for box in boxes[1:-1]:
                box.y = int(current_y + gap)
                current_y = box.y + box.height

        self.refresh_canvas()
        self.mark_unsaved()
        self.update_status(f"已均匀分布 {len(self.selected_boxes)} 个框 ✓")

    def unify_size(self, size_type):
        """统一文本框大小"""
        if len(self.selected_boxes) < 2:
            self.update_status("请Ctrl+点击选中至少2个框")
            return

        self.save_state()
        boxes = [self.text_boxes[i] for i in self.selected_boxes]

        base_box = boxes[0]

        for box in boxes[1:]:
            if size_type in ["width", "both"]:
                box.width = base_box.width
            if size_type in ["height", "both"]:
                box.height = base_box.height

        self.refresh_canvas()
        self.mark_unsaved()
        self.update_status(f"已统一 {len(self.selected_boxes)} 个框的尺寸 ✓")

    def align_to_canvas(self, align_type):
        """对齐到画布中心"""
        if not self.selected_boxes or not self.original_image:
            return

        self.save_state()

        img_w, img_h = self.original_image.size
        center_x = img_w // 2
        center_y = img_h // 2

        for idx in self.selected_boxes:
            box = self.text_boxes[idx]

            if align_type == "h":
                box.x = center_x - box.width // 2
            elif align_type == "v":
                box.y = center_y - box.height // 2
            elif align_type == "center":
                box.x = center_x - box.width // 2
                box.y = center_y - box.height // 2

        self.refresh_canvas()
        self.mark_unsaved()
        self.update_status(f"已对齐到画布中心 ✓")

    # ==================== 新增功能：自动保存 ====================

    def start_autosave(self):
        """启动自动保存"""
        interval = self.config.get("autosave_interval", 300) * 1000
        self.autosave_timer = self.root.after(interval, self.auto_save)

    def stop_autosave(self):
        """停止自动保存"""
        if self.autosave_timer:
            self.root.after_cancel(self.autosave_timer)
            self.autosave_timer = None

    def auto_save(self):
        """自动保存"""
        if self.has_unsaved_changes and self.pages:
            try:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                autosave_path = os.path.join(self.autosave_dir, f"autosave_{timestamp}.json")

                self.save_current_page()

                pages_data = []
                for page in self.pages:
                    pages_data.append({
                        "original_path": page["original_path"],
                        "original_size": page.get("original_size", page["image"].size),
                        "edit_scale": page.get("edit_scale", 1.0),
                        "bg_path": page.get("bg_path"),
                        "text_boxes": page.get("text_boxes", [])
                    })

                with open(autosave_path, 'w', encoding='utf-8') as f:
                    json.dump({"version": 3, "pages": pages_data, "current_page": self.current_page_index},
                             f, ensure_ascii=False, indent=2)

                self.cleanup_autosave_files()
                print(f"自动保存完成: {autosave_path}")

            except Exception as e:
                print(f"自动保存失败: {e}")

        self.start_autosave()

    def cleanup_autosave_files(self):
        """清理旧的自动保存文件"""
        try:
            autosave_files = [f for f in os.listdir(self.autosave_dir) if f.startswith("autosave_")]
            autosave_files.sort(reverse=True)

            for old_file in autosave_files[10:]:
                try:
                    os.remove(os.path.join(self.autosave_dir, old_file))
                except:
                    pass
        except:
            pass

    def mark_unsaved(self):
        """标记有未保存的更改"""
        self.has_unsaved_changes = True
        if hasattr(self, 'autosave_indicator'):
            self.autosave_indicator.config(fg="#FFC107")

    def mark_saved(self):
        """标记已保存"""
        self.has_unsaved_changes = False
        if hasattr(self, 'autosave_indicator'):
            self.autosave_indicator.config(fg="#4CAF50")

    def on_closing(self):
        """窗口关闭事件"""
        if self.has_unsaved_changes:
            result = messagebox.askyesnocancel(
                "未保存的更改",
                "是否保存当前项目？\n\n是 - 保存并退出\n否 - 不保存退出\n取消 - 返回编辑"
            )
            if result is None:
                return
            elif result:
                self.save_project()

        self.stop_autosave()
        self.root.destroy()

    # ==================== 新增功能：PDF导入 ====================

    def import_pdf(self):
        """导入PDF文件 - 使用PyMuPDF，简单快速"""
        if not PDF_SUPPORT:
            messagebox.showerror("需要安装库",
                "PDF转图片需要安装 PyMuPDF\n\n"
                "请运行以下命令:\n"
                "pip install PyMuPDF\n\n"
                "或者:\n"
                "1. 使用在线工具将PDF转为图片\n"
                "2. 然后用'导入图片'功能导入")
            return

        file_path = filedialog.askopenfilename(
            title="选择PDF文件",
            filetypes=[("PDF文件", "*.pdf")]
        )
        if not file_path:
            return

        self.update_status("正在转换PDF...")

        def convert_pdf():
            try:
                self.root.after(0, lambda: self.update_status("正在解析PDF..."))

                # 打开PDF
                doc = fitz.open(file_path)
                page_count = len(doc)

                if page_count == 0:
                    self.root.after(0, lambda: messagebox.showerror("错误", "PDF文件为空"))
                    doc.close()
                    return

                # 询问是否清空现有页面
                if self.pages:
                    result = messagebox.askyesnocancel(
                        "提示",
                        f"PDF共 {page_count} 页。\n\n是否清空现有页面？\n\n"
                        "是 - 清空后导入\n否 - 追加到现有页面\n取消 - 取消导入"
                    )
                    if result is None:
                        self.root.after(0, lambda: self.update_status("已取消"))
                        doc.close()
                        return
                    elif result:
                        self.root.after(0, lambda: setattr(self, 'pages', []))

                # 创建临时目录
                temp_dir = os.path.join(get_base_dir(), "temp_pdf_imports")
                os.makedirs(temp_dir, exist_ok=True)

                start_index = len(self.pages)

                # 转换每一页
                for page_num in range(page_count):
                    self.root.after(0, lambda idx=page_num+1, total=page_count:
                        self.update_status(f"正在转换第 {idx}/{total} 页..."))

                    # 获取页面
                    page = doc[page_num]

                    # 转换为图片（200 DPI高质量）
                    zoom = 200 / 72  # PDF默认72 DPI，提升到200 DPI
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)

                    # 保存为PNG
                    pdf_basename = os.path.splitext(os.path.basename(file_path))[0]
                    temp_path = os.path.join(temp_dir, f"{pdf_basename}_page_{page_num+1:03d}.png")
                    pix.save(temp_path)

                    # 转换为PIL Image
                    img_data = pix.tobytes("png")
                    from io import BytesIO
                    img = Image.open(BytesIO(img_data))

                    # 添加到页面
                    original_size = img.size
                    edit_img, edit_scale = self._resize_image_for_edit(img)

                    page_data = {
                        "original_path": temp_path,
                        "original_size": original_size,
                        "edit_scale": edit_scale,
                        "bg_path": None,
                        "image": edit_img,
                        "text_boxes": []
                    }
                    self.pages.append(page_data)

                # 关闭PDF
                doc.close()

                # 更新界面
                self.root.after(0, lambda: setattr(self, 'current_page_index', start_index))
                self.root.after(0, self.load_current_page)
                self.root.after(0, self.update_page_label)
                self.root.after(0, self.update_thumbnails)
                self.root.after(0, lambda: self.placeholder_label.place_forget())
                self.root.after(0, lambda: self.update_status(f"PDF转换成功！共 {page_count} 页"))
                self.root.after(0, lambda: messagebox.showinfo("成功",
                    f"PDF转换成功！\n\n"
                    f"共转换 {page_count} 页\n"
                    f"图片保存在：{temp_dir}\n\n"
                    f"现在可以进行OCR识别了"))

            except Exception as e:
                import traceback
                error_msg = traceback.format_exc()
                print(f"PDF转换失败:\n{error_msg}")
                self.root.after(0, lambda: messagebox.showerror("错误",
                    f"PDF转换失败:\n\n{str(e)}\n\n"
                    f"建议:\n"
                    f"1. 检查PDF文件是否损坏\n"
                    f"2. 或使用在线工具转换后导入图片"))
                self.root.after(0, lambda: self.update_status("PDF转换失败"))

        threading.Thread(target=convert_pdf, daemon=True).start()

    # ==================== 新增功能：PDF导出 ====================

    def export_as_pdf(self):
        """导出为PDF文件"""
        if not self.pages:
            self.update_status("没有可导出的内容")
            messagebox.showwarning("提示", "请先导入图片或PDF")
            return

        save_path = filedialog.asksaveasfilename(
            defaultextension=".pdf",
            filetypes=[("PDF文件", "*.pdf")],
            initialfile="output.pdf"
        )
        if not save_path:
            return

        self.update_status("正在生成PDF...")

        def export_pdf():
            try:
                self.root.after(0, self.save_current_page)

                pdf_images = []

                for page_idx, page in enumerate(self.pages):
                    self.root.after(0, lambda idx=page_idx+1, total=len(self.pages):
                        self.update_status(f"正在渲染第 {idx}/{total} 页..."))

                    if page.get("bg_path") and os.path.exists(page["bg_path"]):
                        bg_image = Image.open(page["bg_path"])
                    else:
                        bg_image = page["image"].copy()

                    preview_img = bg_image.copy()
                    if preview_img.mode != "RGB":
                        preview_img = preview_img.convert("RGB")

                    draw = ImageDraw.Draw(preview_img)

                    for box_data in page.get("text_boxes", []):
                        if not box_data.get("text"):
                            continue

                        try:
                            pixel_font_size = int(box_data.get("font_size", 16) * 96 / 72)
                            font_path = self._get_font_path(box_data.get("font_name", "微软雅黑"))

                            if font_path and os.path.exists(font_path):
                                font = ImageFont.truetype(font_path, pixel_font_size)
                            else:
                                font = ImageFont.load_default()

                            color_hex = box_data.get("font_color", "#000000").lstrip('#')
                            r = int(color_hex[0:2], 16)
                            g = int(color_hex[2:4], 16)
                            b = int(color_hex[4:6], 16)

                            x, y = box_data["x"], box_data["y"]
                            w, h = box_data["width"], box_data["height"]

                            try:
                                bbox = draw.textbbox((0, 0), box_data["text"], font=font)
                                text_width = bbox[2] - bbox[0]
                                text_height = bbox[3] - bbox[1]
                            except:
                                text_width = len(box_data["text"]) * pixel_font_size * 0.6
                                text_height = pixel_font_size

                            align = box_data.get("align", "left")
                            if align == "center":
                                text_x = x + (w - text_width) // 2
                            elif align == "right":
                                text_x = x + w - text_width - 3
                            else:
                                text_x = x + 3

                            text_y = y + (h - text_height) // 2
                            draw.text((text_x, text_y), box_data["text"], font=font, fill=(r, g, b))

                        except Exception as e:
                            print(f"绘制文字失败 (页{page_idx+1}): {e}")
                            continue

                    pdf_images.append(preview_img)

                if pdf_images:
                    self.root.after(0, lambda: self.update_status("正在保存PDF文件..."))
                    pdf_images[0].save(
                        save_path,
                        "PDF",
                        save_all=True,
                        append_images=pdf_images[1:],
                        resolution=100.0
                    )

                    self.root.after(0, lambda: messagebox.showinfo("成功",
                        f"PDF导出成功！\n\n"
                        f"共 {len(self.pages)} 页\n"
                        f"保存位置：\n{save_path}"))
                    self.root.after(0, lambda: self.update_status(f"PDF导出成功！"))

            except Exception as e:
                import traceback
                error_msg = traceback.format_exc()
                print(f"PDF导出失败:\n{error_msg}")
                self.root.after(0, lambda: messagebox.showerror("错误",
                    f"PDF导出失败:\n\n{str(e)}"))
                self.root.after(0, lambda: self.update_status("PDF导出失败"))

        threading.Thread(target=export_pdf, daemon=True).start()

    # ==================== 新增功能：图片导出 ====================

    def export_as_images(self):
        """导出为图片序列"""
        if not self.pages:
            self.update_status("没有可导出的内容")
            messagebox.showwarning("提示", "请先导入图片或PDF")
            return

        folder_path = filedialog.askdirectory(title="选择导出目录")
        if not folder_path:
            return

        self._show_image_format_dialog(folder_path)

    def _show_image_format_dialog(self, folder_path):
        """显示图片格式选择对话框"""
        format_dialog = tk.Toplevel(self.root)
        format_dialog.title("选择图片格式")
        format_dialog.geometry("350x220")
        format_dialog.configure(bg=COLOR_WHITE)
        format_dialog.transient(self.root)
        format_dialog.grab_set()

        format_dialog.update_idletasks()
        x = (format_dialog.winfo_screenwidth() - 350) // 2
        y = (format_dialog.winfo_screenheight() - 220) // 2
        format_dialog.geometry(f"+{x}+{y}")

        title_frame = tk.Frame(format_dialog, bg=COLOR_THEME, height=40)
        title_frame.pack(fill=tk.X)
        title_frame.pack_propagate(False)
        tk.Label(title_frame, text="  选择图片格式", bg=COLOR_THEME, fg="white",
                font=(FONT_FAMILY, 11, "bold")).pack(side=tk.LEFT, pady=8)

        content = tk.Frame(format_dialog, bg=COLOR_WHITE, padx=20, pady=15)
        content.pack(fill=tk.BOTH, expand=True)

        format_var = tk.StringVar(value="PNG")
        quality_var = tk.IntVar(value=95)

        format_frame = tk.Frame(content, bg=COLOR_WHITE)
        format_frame.pack(pady=10, fill=tk.X)

        tk.Radiobutton(format_frame, text="PNG - 无损压缩，高质量（推荐）",
                      variable=format_var, value="PNG",
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 10)).pack(anchor="w", pady=3)
        tk.Radiobutton(format_frame, text="JPEG - 有损压缩，文件较小",
                      variable=format_var, value="JPEG",
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 10)).pack(anchor="w", pady=3)

        quality_frame = tk.Frame(content, bg=COLOR_WHITE)
        quality_frame.pack(pady=10, fill=tk.X)

        tk.Label(quality_frame, text="JPEG质量 (1-100):", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        quality_scale = tk.Scale(quality_frame, from_=50, to=100, orient=tk.HORIZONTAL,
                                variable=quality_var, bg=COLOR_WHITE, length=150)
        quality_scale.pack(side=tk.LEFT, padx=10)

        btn_frame = tk.Frame(content, bg=COLOR_WHITE)
        btn_frame.pack(pady=15)

        def start_export():
            img_format = format_var.get()
            quality = quality_var.get()
            format_dialog.destroy()
            self._do_export_images(folder_path, img_format, quality)

        tk.Button(btn_frame, text="开始导出", command=start_export,
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 10),
                 padx=20, pady=5, cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=5)

        tk.Button(btn_frame, text="取消", command=format_dialog.destroy,
                 bg=COLOR_GRAY, fg="white", font=(FONT_FAMILY, 10),
                 padx=20, pady=5, cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=5)

    def _do_export_images(self, folder_path, img_format, quality):
        """执行图片导出"""
        self.update_status("正在导出图片...")

        def export_images():
            try:
                self.root.after(0, self.save_current_page)

                for page_idx, page in enumerate(self.pages):
                    self.root.after(0, lambda idx=page_idx+1, total=len(self.pages):
                        self.update_status(f"正在导出第 {idx}/{total} 页..."))

                    if page.get("bg_path") and os.path.exists(page["bg_path"]):
                        bg_image = Image.open(page["bg_path"])
                    else:
                        bg_image = page["image"].copy()

                    preview_img = bg_image.copy()
                    if preview_img.mode not in ["RGB", "RGBA"]:
                        preview_img = preview_img.convert("RGB")

                    draw = ImageDraw.Draw(preview_img)

                    for box_data in page.get("text_boxes", []):
                        if not box_data.get("text"):
                            continue

                        try:
                            pixel_font_size = int(box_data.get("font_size", 16) * 96 / 72)
                            font_path = self._get_font_path(box_data.get("font_name", "微软雅黑"))

                            if font_path and os.path.exists(font_path):
                                font = ImageFont.truetype(font_path, pixel_font_size)
                            else:
                                font = ImageFont.load_default()

                            color_hex = box_data.get("font_color", "#000000").lstrip('#')
                            r = int(color_hex[0:2], 16)
                            g = int(color_hex[2:4], 16)
                            b = int(color_hex[4:6], 16)

                            x, y = box_data["x"], box_data["y"]
                            w, h = box_data["width"], box_data["height"]

                            try:
                                bbox = draw.textbbox((0, 0), box_data["text"], font=font)
                                text_width = bbox[2] - bbox[0]
                                text_height = bbox[3] - bbox[1]
                            except:
                                text_width = len(box_data["text"]) * pixel_font_size * 0.6
                                text_height = pixel_font_size

                            align = box_data.get("align", "left")
                            if align == "center":
                                text_x = x + (w - text_width) // 2
                            elif align == "right":
                                text_x = x + w - text_width - 3
                            else:
                                text_x = x + 3

                            text_y = y + (h - text_height) // 2
                            draw.text((text_x, text_y), box_data["text"], font=font, fill=(r, g, b))

                        except Exception as e:
                            print(f"绘制文字失败 (页{page_idx+1}): {e}")
                            continue

                    ext = ".png" if img_format == "PNG" else ".jpg"
                    save_path = os.path.join(folder_path, f"page_{page_idx+1:03d}{ext}")

                    if img_format == "PNG":
                        if preview_img.mode == "RGBA":
                            preview_img.save(save_path, "PNG")
                        else:
                            preview_img.convert("RGB").save(save_path, "PNG")
                    else:
                        if preview_img.mode == "RGBA":
                            preview_img = preview_img.convert("RGB")
                        preview_img.save(save_path, "JPEG", quality=quality)

                self.root.after(0, lambda: messagebox.showinfo("成功",
                    f"图片导出成功！\n\n"
                    f"共导出 {len(self.pages)} 张图片\n"
                    f"格式：{img_format}\n"
                    f"保存位置：\n{folder_path}"))
                self.root.after(0, lambda: self.update_status(f"图片导出成功！共 {len(self.pages)} 张"))

            except Exception as e:
                import traceback
                error_msg = traceback.format_exc()
                print(f"图片导出失败:\n{error_msg}")
                self.root.after(0, lambda: messagebox.showerror("错误",
                    f"图片导出失败:\n\n{str(e)}"))
                self.root.after(0, lambda: self.update_status("图片导出失败"))

        threading.Thread(target=export_images, daemon=True).start()


if __name__ == "__main__":
    root = tk.Tk()
    app = ModernPPTEditor(root)
    root.mainloop()
